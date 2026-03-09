import os
import re
import json
import hashlib
import threading
import subprocess
import traceback
import zipfile
import shutil
import warnings
from pathlib import Path
from datetime import datetime
from typing import Optional

# PyTorch pin_memory 경고 필터링 (CPU 환경에서 불필요한 경고)
warnings.filterwarnings('ignore', message='.*pin_memory.*', category=UserWarning)

import tkinter as tk
from tkinter import ttk, filedialog, messagebox

from PIL import Image, ImageTk

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from core_store import CoreStore, build_core_artifact_from_doc, build_event_from_artifact, extract_entities_light
from auto_tagging import auto_tag_document, auto_tag_document_chunked, load_auto_tags
from record_builder import RecordBuilderWindow
from chunking import build_chunks_and_reports
from analysis_viewer import open_analysis_window
from ocr_processor import get_ocr_processor
from tagging_viewer import open_tagging_viewer
from collections import Counter

# 그래프 시각화를 위한 선택적 import
try:
    import matplotlib
    matplotlib.use('TkAgg')
    from matplotlib.figure import Figure
    from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
    import matplotlib.pyplot as plt
    import matplotlib.font_manager as fm
    
    # 한글 폰트 설정
    import platform
    system = platform.system()
    if system == 'Windows':
        font_name = 'Malgun Gothic'
    elif system == 'Darwin':
        font_name = 'AppleGothic'
    else:
        font_name = 'NanumGothic'
    
    plt.rcParams['font.family'] = font_name
    plt.rcParams['axes.unicode_minus'] = False
    
    try:
        font_prop = fm.FontProperties(family=font_name)
    except:
        font_prop = None
    
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False
    font_prop = None
except Exception as e:
    print(f"[app] matplotlib 설정 실패: {e}")
    MATPLOTLIB_AVAILABLE = True
    font_prop = None


# =========================
# Config
# =========================
SUPPORTED_DOC_EXTS = {
    ".docx", ".doc",
    ".ppt", ".pptx",
    ".txt", ".md",
    ".pdf",
    ".hwp",
    ".xls", ".xlsx",
}
SUPPORTED_IMG_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"}

DEFAULT_CUSTOMER = "범용"

CATEGORIES = ["오류/장애", "설정/구성", "절차/가이드", "결과/완료", "참고/개념", "기타"]


# =========================
# Utils
# =========================
def now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def sha256_file(path: str) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def stable_id_for_path(path: str) -> str:
    p = Path(path)
    stat = p.stat()
    base = f"{str(p.resolve())}|{stat.st_size}|{int(stat.st_mtime)}"
    return hashlib.sha256(base.encode("utf-8", errors="ignore")).hexdigest()[:16]


def doc_id_for(path: str) -> str:
    return stable_id_for_path(path)


def image_id_for(path: str) -> str:
    return stable_id_for_path(path)


def write_json(path: str, obj: dict):
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(obj, f, ensure_ascii=False, indent=2)


def append_jsonl(path: str, obj: dict):
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    with open(path, "a", encoding="utf-8") as f:
        f.write(json.dumps(obj, ensure_ascii=False) + "\n")


def read_jsonl(path: str) -> list[dict]:
    p = Path(path)
    if not p.exists():
        return []
    out = []
    with open(p, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                out.append(json.loads(line))
            except Exception:
                continue
    return out


def write_jsonl(path: str, records: list[dict]):
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        for r in records:
            f.write(json.dumps(r, ensure_ascii=False) + "\n")


def load_existing_doc_ids(payload_path: str) -> set[str]:
    ids = set()
    for r in read_jsonl(payload_path):
        did = r.get("doc_id")
        if did:
            ids.add(did)
    return ids


def load_existing_image_ids(queue_path: str) -> set[str]:
    ids = set()
    for r in read_jsonl(queue_path):
        iid = r.get("image_id")
        if iid:
            ids.add(iid)
    return ids


def ensure_dict(obj, key, default):
    if key not in obj or obj[key] is None:
        obj[key] = default
    return obj[key]


def guess_language(text: str) -> str:
    if re.search(r"[가-힣]", text):
        return "ko"
    if re.search(r"[A-Za-z]", text):
        return "en"
    return "unknown"


def safe_filename(name: str) -> str:
    return re.sub(r'[\\/:*?"<>|]+', "_", name).strip().strip(".")


def abs_to_file_uri(path: str) -> str:
    p = Path(path).resolve()
    return "file:///" + str(p).replace("\\", "/")


def xml_safe(s: str) -> str:
    """
    python-docx가 허용하지 않는 XML 비호환 문자(널, 제어문자 등) 제거
    - 탭(\t), 개행(\n), 캐리지리턴(\r)은 허용
    """
    if s is None:
        return ""
    if not isinstance(s, str):
        s = str(s)

    s = s.replace("\x00", "")

    out = []
    for ch in s:
        code = ord(ch)
        if code in (0x9, 0xA, 0xD):
            out.append(ch)
        elif 0x20 <= code <= 0xD7FF:
            out.append(ch)
        elif 0xE000 <= code <= 0xFFFD:
            out.append(ch)
        else:
            continue
    return "".join(out)


def clean_extracted_text(extracted: dict, max_chars: int = 12000) -> str:
    if not extracted:
        return ""

    t = extracted.get("type")
    lines: list[str] = []

    if t == "pptx":
        for s in extracted.get("slides", []):
            for x in s.get("texts", []):
                x = (x or "").strip()
                if x:
                    lines.append(x)

    elif t == "docx":
        for x in extracted.get("paragraphs", []):
            x = (x or "").strip()
            if x:
                lines.append(x)

    elif t == "txt":
        for x in extracted.get("lines", []):
            x = (x or "").strip()
            if x:
                lines.append(x)

    else:
        raw = (extracted.get("text") or "").strip()
        if raw:
            lines = [ln.strip() for ln in raw.splitlines() if ln.strip()]

    filtered = []
    for ln in lines:
        if len(ln) <= 1:
            continue
        if re.fullmatch(r"\d{1,3}", ln):
            continue
        filtered.append(ln)

    seen = set()
    uniq = []
    for ln in filtered:
        if ln in seen:
            continue
        seen.add(ln)
        uniq.append(ln)

    text = "\n".join(uniq)
    if len(text) > max_chars:
        text = text[:max_chars] + "\n...\n(이하 생략: 길이 제한)"
    return text


def normalize_customer(s: str) -> str:
    s = (s or "").strip()
    return s if s else DEFAULT_CUSTOMER


def _normalize_one_tag(tag: str) -> str:
    """
    태그 정규화:
    - 앞에 # 제거
    - 공백 → _
    - 소문자
    """
    t = (tag or "").strip()
    if not t:
        return ""
    if t.startswith("#"):
        t = t[1:]
    t = re.sub(r"\s+", "_", t)
    return t.lower()

def normalize_tags(s: str) -> list[str]:
    """
    입력: "db, error,  로그인 실패, #Tomcat"
    출력: ["db", "error", "로그인_실패", "tomcat"]

    - 쉼표(,)로 분리
    - 각 태그는 _normalize_one_tag 규칙 적용
    - 중복 제거 + 정렬
    """
    s = (s or "").strip()
    if not s:
        return []

    parts = [p.strip() for p in s.split(",")]
    tag_set = set()

    for p in parts:
        nt = _normalize_one_tag(p)
        if nt:
            tag_set.add(nt)

    return sorted(tag_set)


def collect_doc_tags_from_images(images: list[dict]) -> list[str]:
    """
    문서에 속한 모든 이미지 메타데이터에서 tags 수집 (중복 제거)
    """
    tag_set = set()
    for im in images:
        if im.get("ignored"):
            continue
        for t in im.get("tags", []) or []:
            nt = _normalize_one_tag(t)
            if nt:
                tag_set.add(nt)
    return sorted(tag_set)


def format_hashtag_line(tags: list[str]) -> str:
    """
    ['db','error','login_fail'] → '#db #error #login_fail'
    """
    if not tags:
        return "(없음)"
    return " ".join(f"#{t}" for t in tags)
# =========================
# Storage merge
# =========================
def upsert_doc_images_metadata(out_root: str, doc_id: str, image_item: dict):
    doc_dir = Path(out_root) / doc_id
    path = doc_dir / "images_metadata.json"

    data = {"doc_id": doc_id, "images": []}
    if path.exists():
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception:
            data = {"doc_id": doc_id, "images": []}

    images = data.get("images", [])
    idx = None
    for i, it in enumerate(images):
        if it.get("image_id") == image_item.get("image_id"):
            idx = i
            break

    if idx is None:
        images.append(image_item)
    else:
        images[idx] = image_item

    data["images"] = images
    write_json(str(path), data)


def upsert_standalone_images_metadata(out_root: str, image_item: dict):
    # 순수 이미지 사전 저장소
    path = Path(out_root) / "standalone_images_metadata.json"

    data = {"images": []}
    if path.exists():
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            data = {"images": []}

    images = data.get("images", [])
    idx = None
    for i, it in enumerate(images):
        if it.get("image_id") == image_item.get("image_id"):
            idx = i
            break

    if idx is None:
        images.append(image_item)
    else:
        images[idx] = image_item

    data["images"] = images
    write_json(str(path), data)


def load_doc_facts(out_root: str, doc_id: str) -> dict | None:
    p = Path(out_root) / doc_id / "document_facts.json"
    if not p.exists():
        return None
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return None


def load_doc_extracted(out_root: str, doc_id: str) -> dict | None:
    p = Path(out_root) / doc_id / "extracted.json"
    if not p.exists():
        return None
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return None


def load_doc_images_metadata(out_root: str, doc_id: str) -> dict | None:
    p = Path(out_root) / doc_id / "images_metadata.json"
    if not p.exists():
        return None
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return None


def load_standalone_images_metadata(out_root: str) -> dict:
    p = Path(out_root) / "standalone_images_metadata.json"
    if not p.exists():
        return {"images": []}
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {"images": []}


# =========================
# Document Metadata Helpers
# =========================
def _datetime_to_iso(dt) -> Optional[str]:
    """datetime 객체를 ISO 형식 문자열로 변환"""
    if dt is None:
        return None
    try:
        if hasattr(dt, "isoformat") and callable(getattr(dt, "isoformat")):
            return dt.isoformat(timespec="seconds")
        return str(dt)
    except Exception:
        return None


def _extract_pdf_metadata(path: str) -> dict:
    """PDF 내부 메타데이터 추출 (제목, 작성일)"""
    result = {"document_title": None, "document_created_at": None}
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        meta = reader.metadata
        if meta:
            if hasattr(meta, "title") and meta.title and str(meta.title).strip():
                result["document_title"] = str(meta.title).strip()
            if hasattr(meta, "creation_date") and meta.creation_date:
                result["document_created_at"] = _datetime_to_iso(meta.creation_date)
    except Exception:
        pass
    return result


def _extract_docx_metadata(path: str) -> dict:
    """DOCX 내부 메타데이터 추출 (제목, 작성일)"""
    result = {"document_title": None, "document_created_at": None}
    try:
        doc = Document(path)
        cp = doc.core_properties
        if cp:
            if hasattr(cp, "title") and cp.title and str(cp.title).strip():
                result["document_title"] = str(cp.title).strip()
            if hasattr(cp, "created") and cp.created:
                result["document_created_at"] = _datetime_to_iso(cp.created)
    except Exception:
        pass
    return result


def _extract_pptx_metadata(path: str) -> dict:
    """PPTX 내부 메타데이터 추출 (제목, 작성일)"""
    result = {"document_title": None, "document_created_at": None}
    try:
        prs = Presentation(path)
        cp = prs.core_properties
        if cp:
            if hasattr(cp, "title") and cp.title and str(cp.title).strip():
                result["document_title"] = str(cp.title).strip()
            if hasattr(cp, "created") and cp.created:
                result["document_created_at"] = _datetime_to_iso(cp.created)
    except Exception:
        pass
    return result


def _extract_xlsx_metadata(path: str) -> dict:
    """XLSX 내부 메타데이터 추출 (제목, 작성일)"""
    result = {"document_title": None, "document_created_at": None}
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        try:
            props = getattr(wb, "properties", None)
            if props:
                if getattr(props, "title", None) and str(props.title or "").strip():
                    result["document_title"] = str(props.title).strip()
                if getattr(props, "created", None):
                    result["document_created_at"] = _datetime_to_iso(props.created)
        finally:
            wb.close()
    except Exception:
        pass
    return result


def extract_document_metadata(path: str, doc_type: str) -> dict:
    """문서 유형별 내부 메타데이터 추출 (제목, 작성일)"""
    ext = (doc_type or "").lower()
    if ext == "pdf":
        return _extract_pdf_metadata(path)
    if ext == "docx":
        return _extract_docx_metadata(path)
    if ext == "pptx":
        return _extract_pptx_metadata(path)
    if ext == "xlsx":
        return _extract_xlsx_metadata(path)
    return {"document_title": None, "document_created_at": None}


# =========================
# Extractors
# =========================
def read_txt(path: str) -> dict:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return {"type": "txt", "lines": lines, "text": text}


def read_docx_text(path: str) -> dict:
    doc = Document(path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]

    tables = []
    for t in doc.tables:
        rows = []
        for r in t.rows:
            rows.append([c.text.strip() for c in r.cells])
        tables.append(rows)

    return {
        "type": "docx",
        "paragraphs": paragraphs,
        "tables": tables,
        "text": "\n".join(paragraphs),
    }


def extract_docx_images(path: str, out_dir: str) -> list[dict]:
    results = []
    doc = Document(path)
    rels = doc.part.related_parts

    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)

    idx = 0
    for _, part in rels.items():
        ctype = getattr(part, "content_type", "")
        if not ctype.startswith("image/"):
            continue

        ext = ctype.split("/")[-1].lower()
        if ext == "jpeg":
            ext = "jpg"

        blob = part.blob
        idx += 1
        fname = f"docx_img_{idx:03d}.{ext}"
        fpath = out_dir_p / fname
        with open(fpath, "wb") as f:
            f.write(blob)

        results.append({
            "image_path": str(fpath.resolve()),
            "origin": "docx_embed",
            "note": ctype
        })
    return results


def read_pptx_text(path: str) -> dict:
    prs = Presentation(path)

    slides = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        slide_tables = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if txt:
                    slide_texts.append(txt)

            if getattr(shape, "has_table", False):
                table = shape.table
                rows = []
                for r in range(len(table.rows)):
                    rows.append([table.cell(r, c).text.strip() for c in range(len(table.columns))])
                slide_tables.append(rows)

        slides.append({"index": i, "texts": slide_texts, "tables": slide_tables})

    return {
        "type": "pptx",
        "slides": slides,
        "text": "\n".join("\n".join(s["texts"]) for s in slides if s["texts"]),
    }


def extract_pptx_images(path: str, out_dir: str) -> list[dict]:
    results = []
    prs = Presentation(path)

    out_dir_p = Path(out_dir)
    out_dir_p.mkdir(parents=True, exist_ok=True)

    idx = 0
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    blob = img.blob
                    ext = (img.ext or "png").lower()
                    if ext == "jpeg":
                        ext = "jpg"
                    idx += 1
                    fname = f"pptx_s{si+1:02d}_img_{idx:03d}.{ext}"
                    fpath = out_dir_p / fname
                    with open(fpath, "wb") as f:
                        f.write(blob)
                    results.append({
                        "image_path": str(fpath.resolve()),
                        "origin": "pptx_embed",
                        "slide_index": si,
                        "note": f"ext={ext}"
                    })
                except Exception:
                    continue

    return results


def convert_ppt_to_pptx_windows(ppt_path: str, out_dir: str, soffice_path: str) -> str:
    ppt_path = str(Path(ppt_path).resolve())
    out_dir = str(Path(out_dir).resolve())
    soffice_path = str(Path(soffice_path).resolve())

    cmd = [
        soffice_path,
        "--headless",
        "--nologo",
        "--norestore",
        "--convert-to", "pptx",
        "--outdir", out_dir,
        ppt_path,
    ]
    subprocess.run(cmd, check=True, capture_output=True, text=True)
    return str(Path(out_dir) / (Path(ppt_path).stem + ".pptx"))
def convert_with_soffice_windows(src_path: str, out_dir: str, soffice_path: str, to_ext: str) -> str:
    src_path = str(Path(src_path).resolve())
    out_dir = str(Path(out_dir).resolve())
    soffice_path = str(Path(soffice_path).resolve())
    Path(out_dir).mkdir(parents=True, exist_ok=True)

    cmd = [
        soffice_path,
        "--headless",
        "--nologo",
        "--norestore",
        "--convert-to", to_ext,
        "--outdir", out_dir,
        src_path,
    ]
    subprocess.run(cmd, check=True, capture_output=True, text=True)
    return str(Path(out_dir) / (Path(src_path).stem + f".{to_ext}"))


def read_pdf_text(
    path: str,
    enable_donut_fallback: bool = True,
    min_text_chars: int = 200,
    min_chars_per_page: int = 50
) -> dict:
    """
    PDF 텍스트 추출 (pypdf 사용, 실패 시 Donut 적용)
    
    Args:
        path: PDF 파일 경로
        enable_donut_fallback: 텍스트 추출 실패 시 Donut 사용 여부
        min_text_chars: 텍스트 레이어 존재 여부를 판단하는 최소 총 문자수
        min_chars_per_page: 페이지당 최소 문자수 기준
    """
    try:
        from pypdf import PdfReader
    except Exception:
        raise ValueError("PDF 처리를 위해 pypdf 필요: pip install pypdf")

    try:
        reader = PdfReader(path)
        page_list = reader.pages
        pages = len(page_list)
        # 100페이지 이상: 페이지 단위 스트리밍 (메모리 절약)
        if pages >= 100:
            texts = []
            for i, p in enumerate(page_list):
                t = (p.extract_text() or "").strip()
                if t:
                    texts.append(t)
                if (i + 1) % 20 == 0:
                    print(f"[PDF] 스트리밍 추출: {i + 1}/{pages} 페이지")
            extracted_text = "\n".join(texts)
        else:
            texts = []
            for p in page_list:
                t = (p.extract_text() or "").strip()
                if t:
                    texts.append(t)
            extracted_text = "\n".join(texts)
        char_count = len(extracted_text)
        chars_per_page = (char_count / pages) if pages > 0 else 0
        has_text_layer = (char_count >= min_text_chars) or (pages > 0 and chars_per_page >= min_chars_per_page)
        
        # 텍스트가 비어있고 Donut fallback이 활성화된 경우
        if not extracted_text.strip() and enable_donut_fallback:
            print(f"[PDF] pypdf로 텍스트 추출 실패 (빈 텍스트), Donut 적용 시도...")
            try:
                from donut_processor import extract_text_with_donut
                donut_result = extract_text_with_donut(path, max_pages=5)
                if donut_result.get("success"):
                    extracted_text = donut_result.get("text", "")
                    char_count = len(extracted_text)
                    has_text_layer = bool(extracted_text.strip())
                    print(f"[PDF] Donut으로 텍스트 추출 성공: {len(extracted_text)} 문자")
                else:
                    print(f"[PDF] Donut 추출 실패: {donut_result.get('error', 'unknown')}")
            except Exception as e:
                print(f"[PDF] Donut 적용 실패: {e}")
        
        return {
            "type": "pdf",
            "pages": pages,
            "text": extracted_text,
            "char_count": char_count,
            "chars_per_page": chars_per_page,
            "has_text_layer": has_text_layer,
            "extraction_method": "donut" if not extracted_text.strip() and enable_donut_fallback else "pypdf"
        }
    except Exception as e:
        # pypdf 추출 실패 시 Donut 시도
        if enable_donut_fallback:
            print(f"[PDF] pypdf 추출 실패: {e}, Donut 적용 시도...")
            try:
                from donut_processor import extract_text_with_donut
                donut_result = extract_text_with_donut(path, max_pages=5)
                if donut_result.get("success"):
                    return {
                        "type": "pdf",
                        "pages": donut_result.get("pages_processed", 0),
                        "text": donut_result.get("text", ""),
                        "extraction_method": "donut"
                    }
                else:
                    print(f"[PDF] Donut 추출 실패: {donut_result.get('error', 'unknown')}")
            except Exception as donut_error:
                print(f"[PDF] Donut 적용 실패: {donut_error}")
        
        # 모든 방법 실패 시 예외 재발생
        raise ValueError(f"PDF 텍스트 추출 실패: {e}")


def read_xlsx_text(path: str) -> dict:
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True, read_only=True)
    lines = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if vals:
                lines.append(" | ".join(vals))

    return {
        "type": "xlsx",
        "text": "\n".join(lines)
    }


def extract_any_document(path: str, tmp_dir: str, soffice_path: str | None, use_donut_fallback: bool = False) -> dict:
    p = Path(path)
    ext = p.suffix.lower()

    # ---- plain text ----
    if ext in (".txt", ".md"):
        return read_txt(str(p))

    # ---- DOCX ----
    if ext == ".docx":
        if not zipfile.is_zipfile(str(p)):
            raise ValueError(
                "Invalid DOCX: 파일이 docx(zip) 형식이 아닙니다. "
                "Word에서 '다른 이름으로 저장(.docx)'으로 재저장 필요"
            )
        return read_docx_text(str(p))

    # ---- PPTX ----
    if ext == ".pptx":
        return read_pptx_text(str(p))

    # ---- PPT → PPTX ----
    if ext == ".ppt":
        if not soffice_path:
            raise ValueError("ppt 처리를 위해 soffice.exe 경로가 필요합니다.")
        pptx_path = convert_ppt_to_pptx_windows(str(p), tmp_dir, soffice_path)
        return read_pptx_text(pptx_path)

    # =========================
    # 🔽 추가 지원 확장자들
    # =========================

    # ---- PDF ----
    if ext == ".pdf":
        return read_pdf_text(str(p), enable_donut_fallback=use_donut_fallback)

    # ---- DOC → DOCX ----
    if ext == ".doc":
        if not soffice_path:
            raise ValueError("doc 처리를 위해 soffice.exe 경로가 필요합니다.")
        docx_path = convert_with_soffice_windows(str(p), tmp_dir, soffice_path, "docx")
        return read_docx_text(docx_path)

    # ---- XLSX ----
    if ext == ".xlsx":
        return read_xlsx_text(str(p))

    # ---- XLS → XLSX ----
    if ext == ".xls":
        if not soffice_path:
            raise ValueError("xls 처리를 위해 soffice.exe 경로가 필요합니다.")
        xlsx_path = convert_with_soffice_windows(str(p), tmp_dir, soffice_path, "xlsx")
        return read_xlsx_text(xlsx_path)

    # ---- HWP → DOCX ----
    if ext == ".hwp":
        if not soffice_path:
            raise ValueError("hwp 처리를 위해 soffice.exe 경로가 필요합니다.")
        docx_path = convert_with_soffice_windows(str(p), tmp_dir, soffice_path, "docx")
        return read_docx_text(docx_path)

    # ---- unsupported ----
    raise ValueError(f"지원하지 않는 문서 형식: {ext}")
    


# =========================
# Scanning
# =========================
def scan_input(input_dir: str) -> tuple[list[str], list[str]]:
    docs = []
    imgs = []
    for root, _, files in os.walk(input_dir):
        for fn in files:
            ext = Path(fn).suffix.lower()
            full = str(Path(root) / fn)
            if ext in SUPPORTED_DOC_EXTS:
                docs.append(full)
            elif ext in SUPPORTED_IMG_EXTS:
                imgs.append(full)
    return docs, imgs


# =========================
# Pipeline
# =========================
def process_one_doc(file_path: str, out_root: str, tmp_dir: str, soffice_path: str | None, inference_scope: str, topic_sentence: Optional[str] = None, use_semantic: bool = False, enable_ocr: bool = True, max_ocr_images: int = 5) -> dict:
    p = Path(file_path)
    doc_id = doc_id_for(file_path)

    # 문서 텍스트 추출 (예외 처리 강화)
    # PDF의 경우 텍스트 레이어가 부족하면 Donut OCR-free를 자동 적용
    use_donut_fallback = enable_ocr and p.suffix.lower() == ".pdf"
    
    try:
        print(f"[추출] 문서 텍스트 추출 시작: {p.name}")
        extracted = extract_any_document(file_path, tmp_dir, soffice_path, use_donut_fallback=use_donut_fallback)
        text = extracted.get("text", "")
        print(f"[추출] 텍스트 추출 완료: {len(text)} 문자")
        
        # PDF에서 텍스트 레이어가 부족하면 OCR 적용 (Donut 우선, 실패 시 EasyOCR)
        if p.suffix.lower() == ".pdf" and enable_ocr:
            char_count = len(text or "")
            pages = extracted.get("pages", 0) or 0
            chars_per_page = (char_count / pages) if pages > 0 else 0
            has_text_layer = extracted.get("has_text_layer")
            
            text_layer_too_low = (
                (has_text_layer is False) or
                (char_count < 200) or
                (pages > 0 and chars_per_page < 50)
            )
            
            if (not text or len(text.strip()) == 0 or text_layer_too_low):
                print(f"[추출] ⚠️ PDF 텍스트 레이어 부족 감지 -> OCR 적용")
                
                # 1단계: Donut 시도
                donut_success = False
                try:
                    from donut_processor import extract_text_with_donut
                    donut_result = extract_text_with_donut(file_path, max_pages=5)
                    if donut_result.get("success"):
                        donut_text = donut_result.get("text", "")
                        if donut_text.strip():
                            text = donut_text
                            extracted["text"] = text
                            extracted["extraction_method"] = "donut"
                            extracted["pages"] = donut_result.get("pages_processed", pages)
                            print(f"[추출] ✅ Donut으로 텍스트 추출 성공: {len(text)} 문자")
                            donut_success = True
                        else:
                            print(f"[추출] ⚠️ Donut으로도 텍스트 추출 실패 (빈 결과)")
                    else:
                        print(f"[추출] ⚠️ Donut 추출 실패: {donut_result.get('error', 'unknown')}")
                except Exception as donut_error:
                    import traceback
                    print(f"[추출] ⚠️ Donut 적용 실패: {donut_error}")
                    traceback.print_exc()
                
                # 2단계: Donut 실패 시 EasyOCR로 PDF 페이지 OCR 처리
                if not donut_success:
                    print(f"[추출] EasyOCR로 PDF 페이지 OCR 처리 시도...")
                    try:
                        # PDF를 이미지로 변환
                        try:
                            from pdf2image import convert_from_path
                            import subprocess
                            
                            # Poppler 경로 자동 찾기
                            poppler_path = None
                            
                            # 1. PATH에서 찾기
                            try:
                                result = subprocess.run(['where', 'pdftoppm'], 
                                                      capture_output=True, 
                                                      text=True, 
                                                      timeout=2)
                                if result.returncode == 0:
                                    poppler_exe = result.stdout.strip().split('\n')[0]
                                    poppler_path = str(Path(poppler_exe).parent)
                            except:
                                pass
                            
                            # 2. 일반적인 위치에서 찾기
                            if not poppler_path:
                                common_paths = [
                                    Path("C:/poppler/bin"),
                                    Path("C:/poppler/Library/bin"),  # 일부 버전은 Library/bin에 있음
                                    Path.home() / "poppler" / "bin",
                                    Path.home() / "poppler-23.11.0" / "bin",
                                    Path("C:/Program Files/poppler/bin"),
                                ]
                                for path in common_paths:
                                    if path.exists() and (path / "pdftoppm.exe").exists():
                                        poppler_path = str(path)
                                        print(f"[OCR] Poppler 경로 발견: {poppler_path}")
                                        break
                                
                                # 3. C:\poppler 안에서 재귀적으로 찾기
                                if not poppler_path:
                                    poppler_root = Path("C:/poppler")
                                    if poppler_root.exists():
                                        for exe_file in poppler_root.rglob("pdftoppm.exe"):
                                            poppler_path = str(exe_file.parent)
                                            print(f"[OCR] Poppler 경로 발견 (재귀 검색): {poppler_path}")
                                            break
                            
                            # convert_from_path 호출 (poppler_path 지정)
                            if poppler_path:
                                pdf_images = convert_from_path(
                                    str(file_path),
                                    dpi=200,
                                    first_page=1,
                                    last_page=min(pages if pages > 0 else 5, 5),
                                    poppler_path=poppler_path
                                )
                            else:
                                # poppler_path 없이 시도 (PATH에 있으면 작동)
                                pdf_images = convert_from_path(
                                    str(file_path),
                                    dpi=200,
                                    first_page=1,
                                    last_page=min(pages if pages > 0 else 5, 5)
                                )
                            
                            print(f"[OCR] PDF 페이지 이미지 변환 완료: {len(pdf_images)}개")
                        except ImportError:
                            print(f"[OCR] ⚠️ pdf2image가 설치되지 않았습니다.")
                            print(f"[OCR] 설치: pip install pdf2image")
                            print(f"[OCR] Windows의 경우 Poppler도 필요합니다:")
                            print(f"[OCR]   - 다운로드: https://github.com/oschwartz10612/poppler-windows/releases")
                            print(f"[OCR]   - 압축 해제 후 bin 폴더를 PATH에 추가하거나")
                            print(f"[OCR]   - convert_from_path(poppler_path='C:/path/to/poppler/bin') 사용")
                            pdf_images = []
                        except Exception as pdf_error:
                            error_msg = str(pdf_error)
                            if "poppler" in error_msg.lower() or "PDFInfoNotInstalledError" in str(type(pdf_error).__name__):
                                print(f"[OCR] ⚠️ Poppler가 설치되지 않았거나 PATH에 없습니다.")
                                print(f"[OCR] Windows 설치 방법:")
                                print(f"[OCR]   1. 다운로드: https://github.com/oschwartz10612/poppler-windows/releases")
                                print(f"[OCR]   2. 압축 해제 (예: C:\\poppler)")
                                print(f"[OCR]   3. 환경 변수 PATH에 C:\\poppler\\bin 추가")
                                print(f"[OCR]   4. 또는 코드에서 poppler_path 지정:")
                                print(f"[OCR]      convert_from_path(..., poppler_path='C:/poppler/bin')")
                            else:
                                print(f"[OCR] PDF 이미지 변환 실패: {pdf_error}")
                            pdf_images = []
                        
                        if pdf_images:
                            # OCR 프로세서 초기화
                            ocr_processor = get_ocr_processor(languages=['ko', 'en'], gpu=False)
                            
                            # 각 페이지 OCR 처리
                            ocr_texts = []
                            import tempfile
                            for idx, img in enumerate(pdf_images, start=1):
                                print(f"[OCR] PDF 페이지 {idx}/{len(pdf_images)} OCR 처리 중...")
                                try:
                                    # 임시 이미지 파일로 저장
                                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_file:
                                        tmp_path = tmp_file.name
                                        img.save(tmp_path, "PNG")
                                    
                                    # OCR 처리
                                    ocr_result = ocr_processor.extract_text(tmp_path)
                                    if ocr_result.get("success") and ocr_result.get("text"):
                                        page_text = ocr_result.get("text", "").strip()
                                        if page_text:
                                            ocr_texts.append(f"=== 페이지 {idx} ===\n{page_text}")
                                    
                                    # 임시 파일 삭제
                                    Path(tmp_path).unlink(missing_ok=True)
                                except Exception as page_error:
                                    print(f"[OCR] 페이지 {idx} OCR 실패: {page_error}")
                                    Path(tmp_path).unlink(missing_ok=True)
                            
                            # OCR 결과 결합
                            if ocr_texts:
                                ocr_combined_text = "\n\n".join(ocr_texts)
                                if ocr_combined_text.strip():
                                    text = ocr_combined_text
                                    extracted["text"] = text
                                    extracted["extraction_method"] = "easyocr_pdf"
                                    print(f"[추출] ✅ EasyOCR로 PDF 텍스트 추출 성공: {len(text)} 문자 ({len(ocr_texts)}페이지)")
                                else:
                                    print(f"[추출] ⚠️ EasyOCR 결과가 비어있습니다")
                            else:
                                print(f"[추출] ⚠️ EasyOCR로 텍스트 추출 실패 (모든 페이지 실패)")
                    except Exception as ocr_error:
                        import traceback
                        print(f"[추출] ⚠️ EasyOCR PDF 처리 실패: {ocr_error}")
                        traceback.print_exc()
        
        # 텍스트가 비어있을 경우 경고
        if not text or len(text.strip()) == 0:
            print(f"[추출] ⚠️ 경고: 추출된 텍스트가 비어있습니다. 파일: {p.name}")
            print(f"[추출] 문서 타입: {extracted.get('type', 'unknown')}")
    except Exception as e:
        import traceback
        print(f"[추출] ❌ 문서 텍스트 추출 실패: {p.name}")
        print(f"[추출] 오류: {e}")
        traceback.print_exc()
        
        # PDF이고 OCR 활성화된 경우 Donut 재시도
        if p.suffix.lower() == ".pdf" and enable_ocr:
            print(f"[추출] Donut OCR-free 적용 시도...")
            try:
                from donut_processor import extract_text_with_donut
                donut_result = extract_text_with_donut(file_path, max_pages=5)
                if donut_result.get("success"):
                    text = donut_result.get("text", "")
                    extracted = {
                        "type": "pdf",
                        "pages": donut_result.get("pages_processed", 0),
                        "text": text,
                        "extraction_method": "donut"
                    }
                    print(f"[추출] ✅ Donut으로 텍스트 추출 성공: {len(text)} 문자")
                else:
                    print(f"[추출] Donut 추출 실패: {donut_result.get('error', 'unknown')}")
                    extracted = {"type": p.suffix.lower().lstrip("."), "text": ""}
                    text = ""
            except Exception as donut_error:
                print(f"[추출] Donut 적용 실패: {donut_error}")
                extracted = {"type": p.suffix.lower().lstrip("."), "text": ""}
                text = ""
        else:
            # 추출 실패 시 빈 텍스트로 계속 진행 (태깅은 실패하지만 프로세스는 계속)
            extracted = {"type": p.suffix.lower().lstrip("."), "text": ""}
            text = ""

    stat = p.stat()
    ingested_at = now_iso()
    doc_type = p.suffix.lower().lstrip(".")
    
    # 문서 내부 메타데이터 추출 (제목, 작성일)
    doc_meta = extract_document_metadata(str(p.resolve()), doc_type)
    document_title = doc_meta.get("document_title")
    document_created_at = doc_meta.get("document_created_at")
    display_title = (document_title or "").strip() or p.stem
    
    facts = {
        "doc_id": doc_id,
        "source_path": str(p.resolve()),
        "doc_type": doc_type,
        "file_name": p.name,
        "title": display_title,
        "document_title": document_title,
        "document_created_at": document_created_at,
        "ingested_at": ingested_at,
        "size_bytes": stat.st_size,
        "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds"),
        "sha256": sha256_file(file_path),
        "text_length": len(text),
        "language": guess_language(text),
        "inference_scope": inference_scope,
        "processed_at": ingested_at,
    }

    doc_dir = Path(out_root) / doc_id
    write_json(str(doc_dir / "document_facts.json"), facts)
    write_json(str(doc_dir / "extracted.json"), extracted)

    extracted_images = []
    images_dir = doc_dir / "images"

    if p.suffix.lower() == ".docx":
        extracted_images = extract_docx_images(str(p), str(images_dir))
    elif p.suffix.lower() == ".pptx":
        extracted_images = extract_pptx_images(str(p), str(images_dir))
    elif p.suffix.lower() == ".ppt":
        converted = str(Path(tmp_dir) / (p.stem + ".pptx"))
        if Path(converted).exists():
            extracted_images = extract_pptx_images(converted, str(images_dir))
    
    # OCR 처리 (자동, CPU 모드) - 메모리 절약을 위해 선택적 사용
    ocr_processed_count = 0
    if extracted_images and enable_ocr:
        try:
            # OCR은 메모리를 많이 사용하므로, 이미지가 많으면 건너뛰기 옵션
            ocr_processor = get_ocr_processor(languages=['ko', 'en'], gpu=False)
            # 메모리 절약: 최대 지정된 개수만 처리
            images_to_process = extracted_images[:max_ocr_images] if max_ocr_images > 0 else []
            if len(extracted_images) > max_ocr_images and max_ocr_images > 0:
                print(f"[OCR] 메모리 절약: {len(extracted_images)}개 중 {max_ocr_images}개만 처리")
            
            for img_info in images_to_process:
                img_path = img_info.get("image_path")
                if img_path and Path(img_path).exists():
                    try:
                        ocr_result = ocr_processor.extract_text(img_path)
                        if ocr_result.get("success"):
                            # OCR 결과를 이미지 정보에 추가
                            img_info["ocr"] = {
                                "text": ocr_result.get("text", ""),
                                "lines": ocr_result.get("lines", []),
                                "confidence": ocr_result.get("confidence", 0.0),
                                "word_count": ocr_result.get("word_count", 0),
                                "line_count": ocr_result.get("line_count", 0),
                                "language": ocr_result.get("language", "unknown"),
                                "processed_at": now_iso()
                            }
                            ocr_processed_count += 1
                    except Exception as e:
                        # 개별 이미지 OCR 실패는 무시 (다음 이미지 계속 처리)
                        import traceback
                        print(f"[OCR] 이미지 OCR 실패 ({Path(img_path).name}): {e}")
                        pass
        except ImportError:
            # EasyOCR이 설치되지 않은 경우
            print("[OCR] EasyOCR이 설치되지 않았습니다. 'pip install easyocr'을 실행하세요.")
            pass
        except Exception as e:
            # OCR 초기화 실패
            import traceback
            print(f"[OCR] OCR 초기화 실패: {e}")
            traceback.print_exc()
            pass

    payload = {
        "doc_id": doc_id,
        "title": display_title,
        "document_title": document_title,
        "document_created_at": document_created_at,
        "source_path": facts["source_path"],
        "ingested_at": ingested_at,
        "doc_type": facts["doc_type"],
        "language": facts["language"],
        "inference_scope": facts["inference_scope"],
        "modified_at": facts["modified_at"],
        "text": text,
        "stats": {"size_bytes": facts["size_bytes"], "text_length": facts["text_length"]},
        "assets": {"extracted_images": extracted_images},
    }
    append_jsonl(str(Path(out_root) / "workspace_payload.jsonl"), payload)
        # ✅ 자동 태깅(문서 텍스트 기반) 생성/저장
    # OCR과 태깅은 완전히 독립적입니다. OCR이 없어도 태깅은 정상 작동합니다.
    try:
        print(f"[태깅] 문서 태깅 시작: {doc_id}")
        print(f"[태깅] 텍스트 길이: {len(text)} 문자")
        print(f"[태깅] 설정 - 의미보정: {use_semantic}, OCR: {enable_ocr}, 최대OCR: {max_ocr_images}")
        
        if not text or len(text.strip()) == 0:
            print(f"[태깅] ⚠️ 경고: 텍스트가 비어있습니다. 태깅을 건너뜁니다.")
        elif len(text) > 100000:
            print(f"[태깅] 대용량 문서({len(text):,}자) - 청크 단위 태깅 적용")
            auto_tags_result = auto_tag_document_chunked(
                out_root=out_root,
                doc_id=doc_id,
                title=facts.get("title", ""),
                text=text or "",
                chunk_size=50000,
                language=facts.get("language", "unknown"),
                top_k=15,
                topic_sentence=topic_sentence,
                min_algorithm_support=1,
                use_semantic_adjustment=use_semantic,
            )
        else:
            auto_tags_result = auto_tag_document(
                out_root=out_root,
                doc_id=doc_id,
                title=facts.get("title", ""),
                text=text or "",
                language=facts.get("language", "unknown"),
                top_k=15,  # 최소 10개 보장을 위해 15개 요청
                topic_sentence=topic_sentence,
                min_algorithm_support=1,  # 최소 지지 수를 1로 낮춰 더 많은 태그 생성
                use_semantic_adjustment=use_semantic  # 메모리 절약을 위해 선택적 사용
            )
            print(f"[태깅] ✅ 문서 태깅 완료: {doc_id}")
            
            # PDF에서 장르가 미결정(unknown)이고 OCR이 활성화된 경우 Donut 재시도
            if (p.suffix.lower() == ".pdf" and enable_ocr and
                auto_tags_result and auto_tags_result.get("genre") == "unknown"):
                print(f"[태깅] ⚠️ PDF 장르 미결정, Donut OCR-free로 재처리 시도...")
                # Donut으로 텍스트 재추출 후 재태깅
                try:
                    from donut_processor import extract_text_with_donut
                    donut_result = extract_text_with_donut(file_path, max_pages=5)
                    if donut_result.get("success"):
                        donut_text = donut_result.get("text", "")
                        if donut_text.strip() and len(donut_text.strip()) > len(text.strip()):
                            print(f"[태깅] Donut 텍스트로 재태깅 시도... (기존: {len(text)}자, Donut: {len(donut_text)}자)")
                            auto_tag_document(
                                out_root=out_root,
                                doc_id=doc_id,
                                title=facts.get("title", ""),
                                text=donut_text,
                                language=facts.get("language", "unknown"),
                                top_k=15,
                                topic_sentence=topic_sentence,
                                min_algorithm_support=1,
                                use_semantic_adjustment=use_semantic
                            )
                            print(f"[태깅] ✅ Donut 텍스트로 재태깅 완료")
                        else:
                            print(f"[태깅] Donut 텍스트가 기존과 같거나 짧아서 재태깅 건너뜀")
                except Exception as donut_error:
                    import traceback
                    print(f"[태깅] Donut 재태깅 실패: {donut_error}")
                    traceback.print_exc()
    except Exception as e:
        import traceback
        print(f"[태깅] ❌ 문서 태깅 실패: {doc_id}")
        print(f"[태깅] 오류: {e}")
        traceback.print_exc()

    # ============================================================
    # ✅ 자동 청킹 처리 (태깅 완료 후 자동 실행)
    # - 128~1024 토큰 범위 (문자 수 기반: 약 128~1024자)
    # - 장르별 최적화된 청킹 전략 적용
    # ============================================================
    try:
        if text and len(text.strip()) > 0:
            print(f"[청킹] 문서 청킹 시작: {doc_id}")
            print(f"[청킹] 텍스트 길이: {len(text)} 문자")
            
            # 청킹 크기 설정 (128~1024 토큰 범위, 문자 수 기반)
            # 오버랩 확대: 문맥 연속성 향상 (target의 약 15~20%)
            target_chars = 512
            min_chars = 128
            max_chars = 1024
            overlap_chars = 96  # 확대: 64→96 (target의 약 19%)
            
            if len(text) < 500:
                target_chars = min(256, len(text) // 2)
                min_chars = min(64, len(text) // 4)
                overlap_chars = 48
            elif len(text) > 10000:
                target_chars = 768
                min_chars = 256
                overlap_chars = 128  # 긴 문서: 오버랩 확대
            
            print(f"[청킹] 청킹 설정: target={target_chars}, min={min_chars}, overlap={overlap_chars}")
            
            build_chunks_and_reports(
                out_root=out_root,
                doc_id=doc_id,
                facts=facts,
                text=text,
                target_chars=target_chars,
                min_chars=min_chars,
                overlap_chars=overlap_chars
            )
            print(f"[청킹] ✅ 문서 청킹 완료: {doc_id}")
        else:
            print(f"[청킹] ⚠️ 경고: 텍스트가 비어있습니다. 청킹을 건너뜁니다.")
    except Exception as e:
        import traceback
        print(f"[청킹] ❌ 문서 청킹 실패: {doc_id}")
        print(f"[청킹] 오류: {e}")
        traceback.print_exc()
        # 청킹 실패해도 프로세스는 계속 진행

    # ============================================================
    # ✅ Core Schema v1 출력(범용화: 다마고치 AI에도 그대로 사용)
    # - artifacts.jsonl
    # - entities.jsonl (light)
    # - events.jsonl
    # - records/<entity_id>.json (light timeline)
    # ============================================================
    try:
        core = CoreStore(out_root)

        # attachments: 문서 내 추출 이미지 링크만 기록(라벨링은 기존 흐름 그대로 유지)
        attachments = []
        for it in extracted_images or []:
            ip = it.get("image_path")
            if ip:
                attachments.append({"kind": "image", "ref": abs_to_file_uri(ip), "note": it.get("origin", "")})

        artifact = build_core_artifact_from_doc(
            source_path=facts["source_path"],
            sha256=facts["sha256"],
            size_bytes=facts["size_bytes"],
            modified_at=facts["modified_at"],
            ingested_at=facts["processed_at"],
            inference_scope=facts["inference_scope"],
            title=facts["title"],
            language=facts["language"],
            text=text,
            attachments=attachments,
        )
        artifact_id = core.append_artifact(artifact)

        # 엔티티 추출(v1 light) + upsert
        ents = extract_entities_light(text=text, title=facts["title"])
        entity_refs = []
        for e in ents:
            er = core.upsert_entity(e.get("type", "other"), e.get("name", ""), artifact_id)
            if er and er.get("entity_id"):
                entity_refs.append({"entity_id": er["entity_id"], "role": "related"})

        # 문서 자체도 하나의 entity(topic)로 연결하고 싶으면(선택):
        # doc_entity = core.upsert_entity("doc_topic", facts["title"], artifact_id)
        # if doc_entity.get("entity_id"):
        #     entity_refs.append({"entity_id": doc_entity["entity_id"], "role": "primary"})

        genre = (artifact.get("classification") or {}).get("genre", "unknown")
        ev_summary = f"[ingest] {facts['title']} ({genre})"

        ev = build_event_from_artifact(
            artifact_id=artifact_id,
            event_time=facts["processed_at"],
            summary=ev_summary,
            entities=entity_refs,
            genre=genre,
            evidence=[{
                "artifact_id": artifact_id,
                "quote": (text.strip().splitlines()[0][:180] if text.strip() else "(empty)"),
                "locator": "line1"
            }]
        )
        event_id = core.append_event(ev)

        # records 업데이트: 엔티티별 타임라인에 이벤트를 넣어줌
        for r in entity_refs:
            eid = r.get("entity_id")
            if not eid:
                continue
            core.apply_event_to_record(
                entity_id=eid,
                event_id=event_id,
                event_time=facts["processed_at"],
                summary=ev_summary,
                artifact_id=artifact_id,
                artifact_title=facts["title"],
                genre=genre
            )

    except Exception:
        # 코어 저장은 실패해도 기존 파이프라인을 절대 깨지 않게
        pass

    return {
        "ok": True,
        "doc_id": doc_id,
        "path": str(p),
        "extracted_images": extracted_images,
    }


def register_image_for_review(
    out_root: str,
    image_path: str,
    origin: str,
    mode: str,  # "doc_embed" or "standalone"
    parent_doc_id: str | None,
    parent_path: str | None,
    location: dict | None,
    existing_image_ids: set[str],
    ocr_info: dict | None = None,  # OCR 정보 (있는 경우)
) -> dict:
    img_p = Path(image_path)
    stat = img_p.stat()
    iid = image_id_for(str(img_p))

    # label: 범용 5문항 스키마
    rec = {
        "image_id": iid,
        "image_path": str(img_p.resolve()),
        "origin": origin,
        "mode": mode,
        "parent_doc_id": parent_doc_id,
        "parent_source_path": parent_path,
        "location": location or {},
        "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds"),
        "sha256": sha256_file(str(img_p)),
        "labeled": False,
        "ignored": False,
        "label": {
            "title": img_p.stem,          # 기본값: 파일명
            "customer": DEFAULT_CUSTOMER, # 기본값
            "category": CATEGORIES[0],
            "description": "",
            "tags": [],
        },
        "status": {"saved_label": False, "merged": False, "ignored": False},
        "created_at": now_iso(),
    }

    # OCR 정보 추가 (있는 경우)
    if ocr_info:
        rec["ocr"] = ocr_info

    if iid in existing_image_ids:
        rec["skipped"] = True
        return rec

    append_jsonl(str(Path(out_root) / "images_queue.jsonl"), rec)
    existing_image_ids.add(iid)
    rec["skipped"] = False
    return rec


# =========================
# Export: Metadata DOCX
# =========================
def generate_metadata_docx_for_doc(out_root: str, doc_id: str) -> tuple[bool, str]:
    facts = load_doc_facts(out_root, doc_id)
    if not facts:
        return False, f"document_facts.json 없음: doc_id={doc_id}"

    extracted = load_doc_extracted(out_root, doc_id) or {}
    clean_text = clean_extracted_text(extracted, max_chars=12000)

    images_meta = load_doc_images_metadata(out_root, doc_id) or {"doc_id": doc_id, "images": []}
    images = images_meta.get("images", [])
    doc_tags = collect_doc_tags_from_images(images)
    
    # OCR 텍스트 수집 (검색용)
    ocr_texts = []
    for im in images:
        ocr_info = im.get("ocr")
        if ocr_info and ocr_info.get("text"):
            ocr_texts.append(ocr_info.get("text", ""))

    # 정제: ignored 제거 + description 없는 이미지 제거
    usable = []
    for im in images:
        if im.get("ignored", False):
            continue
        desc = (im.get("description") or "").strip()
        if not desc:
            continue
        usable.append(im)

    title = facts.get("title") or facts.get("file_name") or doc_id
    base_name = safe_filename(title)
    out_dir = Path(out_root) / "metadata_docs"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{base_name}.docx"

    doc = Document()
    doc.add_heading(xml_safe(f"[정제 메타데이터] {title}"), level=1)

    doc.add_heading("문서 정보", level=2)
    doc.add_paragraph(xml_safe(f"문서 ID: {facts.get('doc_id')}"))
    doc.add_paragraph(xml_safe(f"원본 파일명: {facts.get('file_name')}"))
    doc.add_paragraph(xml_safe(f"원본 경로: {facts.get('source_path')}"))
    doc.add_paragraph(xml_safe(f"문서 유형: {facts.get('doc_type')}"))
    doc.add_paragraph(xml_safe(f"언어: {facts.get('language')}"))
    doc.add_paragraph(xml_safe(f"inference_scope: {facts.get('inference_scope')}"))
    doc.add_paragraph(xml_safe(f"문서 입력일: {facts.get('ingested_at')}"))
    doc.add_paragraph(xml_safe(f"문서 작성일: {facts.get('document_created_at') or '(없음)'}"))
    doc.add_paragraph(xml_safe(f"문서 내부 제목: {facts.get('document_title') or '(없음)'}"))
    doc.add_paragraph(xml_safe(f"수정 시각: {facts.get('modified_at')}"))
    doc.add_paragraph(xml_safe(f"처리 시각: {facts.get('processed_at')}"))
    doc.add_paragraph(xml_safe(f"sha256: {facts.get('sha256')}"))
    doc.add_paragraph(xml_safe(f"문서 태그: {format_hashtag_line(doc_tags)}"))
    
    # OCR 텍스트 추가 (있는 경우)
    if ocr_texts:
        doc.add_heading("OCR 텍스트(이미지에서 추출)", level=2)
        ocr_combined = "\n\n--- 이미지별 OCR 결과 ---\n\n".join(ocr_texts)
        doc.add_paragraph(xml_safe(ocr_combined))

    doc.add_heading("정제 텍스트(추출 기반)", level=2)
    if clean_text.strip():
        doc.add_paragraph(xml_safe(clean_text))
    else:
        doc.add_paragraph("추출된 텍스트가 없거나 정제 결과가 비어있습니다.")

    doc.add_heading("이미지 메타데이터(설명 있는 것만)", level=2)
    if not usable:
        doc.add_paragraph("설명이 입력된 이미지가 없습니다. (무시 처리 또는 설명 미입력)")
    else:
        for idx, im in enumerate(usable, start=1):
            loc = im.get("location") or {}
            loc_str = ""
            if "slide_index" in loc and loc.get("slide_index") is not None:
                loc_str = f"슬라이드 {int(loc['slide_index']) + 1}"
            elif loc:
                loc_str = str(loc)

            img_path = im.get("image_path", "")
            img_uri = abs_to_file_uri(img_path) if img_path else ""

            doc.add_heading(xml_safe(f"IMG-{idx:03d}"), level=3)
            if loc_str:
                doc.add_paragraph(xml_safe(f"위치: {loc_str}"))
            doc.add_paragraph(xml_safe(f"제목: {im.get('title','')}"))
            doc.add_paragraph(xml_safe(f"고객사: {im.get('customer','')}"))
            doc.add_paragraph(xml_safe(f"분류: {im.get('category','')}"))
            doc.add_paragraph(xml_safe(f"설명: {im.get('description','')}"))
            tags = im.get("tags") or []
            doc.add_paragraph(xml_safe(f"태그: {format_hashtag_line(tags)}"))
            
            # OCR 텍스트 추가 (있는 경우)
            ocr_info = im.get("ocr")
            if ocr_info and ocr_info.get("text"):
                ocr_text = ocr_info.get("text", "")
                ocr_conf = ocr_info.get("confidence", 0.0)
                doc.add_paragraph(xml_safe(f"[OCR 텍스트] (신뢰도: {ocr_conf:.2%})\n{ocr_text}"))
            
            doc.add_paragraph(xml_safe(f"참조 이미지 링크: {img_uri}"))
            doc.add_paragraph(xml_safe(f"참조 이미지 경로: {img_path}"))

    doc.add_heading("AI 사용 규칙(권장)", level=2)
    doc.add_paragraph("1) 본 DOCX(정제 메타데이터)가 존재하면 반드시 이를 먼저 참조한다.")
    doc.add_paragraph("2) 이미지 설명(메타데이터)은 사람(문서 제공자)이 정의한 사실로 사용한다. (이미지 자체 추측 해석 금지)")
    doc.add_paragraph("3) 본 문서에 없는 정보에 한해 원본 문서를 추가로 참조한다.")
    doc.add_paragraph("4) 필요 시 답변에 IMG-xxx 및 참조 링크(file://)를 함께 제시한다.")

    doc.save(str(out_path))
    return True, str(out_path)


def generate_image_dictionary_docx(out_root: str) -> tuple[bool, str]:
    data = load_standalone_images_metadata(out_root)
    images = data.get("images", [])

    usable = []
    for im in images:
        if im.get("ignored", False):
            continue
        desc = (im.get("description") or "").strip()
        if not desc:
            continue
        usable.append(im)

    out_dir = Path(out_root) / "metadata_docs"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "_IMAGE_DICTIONARY.docx"

    doc = Document()
    doc.add_heading(xml_safe("[순수 이미지 지식 사전]"), level=1)
    doc.add_paragraph(xml_safe(f"생성 시각: {now_iso()}"))
    doc.add_paragraph(xml_safe("설명(description)이 있는 이미지들만 포함됩니다. (무시/미기입은 제거)"))

    if not usable:
        doc.add_paragraph("현재 포함할 이미지가 없습니다.")
        doc.save(str(out_path))
        return True, str(out_path)

    for idx, im in enumerate(usable, start=1):
        img_path = im.get("image_path", "")
        img_uri = abs_to_file_uri(img_path) if img_path else ""
        doc.add_heading(xml_safe(f"IMG-{idx:03d}"), level=2)
        doc.add_paragraph(xml_safe(f"제목: {im.get('title','')}"))
        doc.add_paragraph(xml_safe(f"고객사: {im.get('customer','')}"))
        doc.add_paragraph(xml_safe(f"분류: {im.get('category','')}"))
        doc.add_paragraph(xml_safe(f"설명: {im.get('description','')}"))
        tags = im.get("tags") or []
        doc.add_paragraph(xml_safe(f"태그: {', '.join(tags) if tags else '(없음)'}"))
        doc.add_paragraph(xml_safe(f"참조 이미지 링크: {img_uri}"))
        doc.add_paragraph(xml_safe(f"참조 이미지 경로: {img_path}"))

    doc.save(str(out_path))
    return True, str(out_path)


def generate_metadata_docx_for_all(out_root: str) -> tuple[int, int, list[str]]:
    out = Path(out_root)
    ok = 0
    fail = 0
    msgs = []

    # 문서별
    for child in out.iterdir():
        if not child.is_dir():
            continue
        if child.name in ("metadata_docs",):
            continue
        facts_path = child / "document_facts.json"
        if facts_path.exists():
            did = child.name
            try:
                success, msg = generate_metadata_docx_for_doc(out_root, did)
                if success:
                    ok += 1
                else:
                    fail += 1
                msgs.append(f"{did}: {msg}")
            except Exception as e:
                fail += 1
                msgs.append(f"{did}: FAIL {e}")

    # 순수 이미지 사전도 함께 생성
    try:
        s_ok, s_path = generate_image_dictionary_docx(out_root)
        msgs.append(f"standalone: {s_path if s_ok else 'FAIL'}")
    except Exception as e:
        msgs.append(f"standalone: FAIL {e}")

    return ok, fail, msgs


# =========================
# UI: Image Review
# =========================
class ImageReviewWindow(tk.Toplevel):
    def __init__(self, master, out_root: str, selected_doc_id: str):
        super().__init__(master)
        self.title("이미지 라벨링(범용 5문항)")
        self.geometry("1080x820")

        self.out_root = out_root
        self.selected_doc_id = selected_doc_id

        self.queue_path = str(Path(out_root) / "images_queue.jsonl")
        self.labels_path = str(Path(out_root) / "images_labels.jsonl")

# 🔥 처음부터 필터된 큐만 로드
        self.records = self._load_filtered_queue()
        self.index = 0

        self._build_ui()
        self._refresh()
    
    def _load_filtered_queue(self):
        all_items = read_jsonl(self.queue_path)
        filtered = []

        # workspace_payload에서 OCR 정보 수집 (이미지 경로 기준, 정규화된 경로 사용)
        ocr_map = {}
        if self.selected_doc_id != "__STANDALONE__":
            payload_path = Path(self.out_root) / "workspace_payload.jsonl"
            if payload_path.exists():
                payloads = read_jsonl(str(payload_path))
                for payload in payloads:
                    if payload.get("doc_id") == self.selected_doc_id:
                        for img_info in payload.get("assets", {}).get("extracted_images", []):
                            img_path = img_info.get("image_path")
                            if img_path and "ocr" in img_info:
                                # 경로를 정규화해서 키로 사용
                                try:
                                    normalized_path = str(Path(img_path).resolve())
                                    ocr_map[normalized_path] = img_info["ocr"]
                                except Exception:
                                    # 경로 정규화 실패 시 원본 경로 사용
                                    ocr_map[img_path] = img_info["ocr"]
                        break

        for img in all_items:
            if self.selected_doc_id == "__STANDALONE__":
                if img.get("mode") == "standalone":
                    filtered.append(img)
            else:
                if img.get("parent_doc_id") == self.selected_doc_id:
                    # OCR 정보가 없으면 workspace_payload에서 가져오기
                    if "ocr" not in img:
                        img_path = img.get("image_path")
                        if img_path:
                            # 정규화된 경로로 검색 시도
                            try:
                                normalized_path = str(Path(img_path).resolve())
                                if normalized_path in ocr_map:
                                    img["ocr"] = ocr_map[normalized_path]
                                elif img_path in ocr_map:
                                    img["ocr"] = ocr_map[img_path]
                            except Exception:
                                # 경로 정규화 실패 시 원본 경로로 검색
                                if img_path in ocr_map:
                                    img["ocr"] = ocr_map[img_path]
                    filtered.append(img)

    # 미라벨 우선 정렬
        filtered.sort(key=lambda x: x.get("labeled", False))
        return filtered



    def _load_queue(self) -> list[dict]:
        recs = read_jsonl(self.queue_path)

        # 구버전 호환: label 구조 보정
        for r in recs:
            ensure_dict(r, "status", {"saved_label": False, "merged": False, "ignored": False})
            ensure_dict(r, "label", {})
            # 기존 meaning/note/tags 구조가 있으면 description/tags로 흡수
            if "meaning" in r["label"] and "description" not in r["label"]:
                r["label"]["description"] = r["label"].get("meaning", "")
            if "tags" not in r["label"]:
                r["label"]["tags"] = []
            if "title" not in r["label"]:
                r["label"]["title"] = Path(r.get("image_path","")).stem
            if "customer" not in r["label"]:
                r["label"]["customer"] = DEFAULT_CUSTOMER
            if "category" not in r["label"]:
                r["label"]["category"] = CATEGORIES[0]
            if "description" not in r["label"]:
                r["label"]["description"] = ""
            if "mode" not in r:
                r["mode"] = "doc_embed" if r.get("parent_doc_id") else "standalone"

        return recs

    def _persist_queue(self):
        write_jsonl(self.queue_path, self.records)

    def _count_remaining(self) -> int:
        cnt = 0
        for r in self.records:
            st = r.get("status", {})
            done = bool(r.get("labeled")) or bool(st.get("saved_label"))
            if not done:
                cnt += 1
        return cnt

    def _build_ui(self):
        pad = {"padx": 10, "pady": 8}

        top = ttk.Frame(self); top.pack(fill="x", **pad)
        self.lbl_pos = ttk.Label(top, text="0/0")
        self.lbl_pos.pack(side="left")

        ttk.Button(top, text="이전", command=self.prev).pack(side="left", padx=6)
        ttk.Button(top, text="다음", command=self.next).pack(side="left", padx=6)
        ttk.Button(top, text="마지막 이미지로", command=self.goto_last).pack(side="left", padx=12)

        ttk.Button(top, text="저장", command=self.save_current).pack(side="left", padx=12)
        ttk.Button(top, text="저장하지 않음(무시)", command=self.mark_ignored_and_next).pack(side="left", padx=6)

        ttk.Button(top, text="완료하기", command=self.finish).pack(side="left", padx=18)
        ttk.Button(top, text="큐 새로고침", command=self.reload).pack(side="right")

        mid = ttk.Frame(self); mid.pack(fill="both", expand=True, **pad)

        left = ttk.Frame(mid); left.pack(side="left", fill="both", expand=True, padx=8)
        self.canvas = tk.Canvas(left, bg="#111", width=680, height=600)
        self.canvas.pack(fill="both", expand=True)

        self.lbl_meta = ttk.Label(left, text="", justify="left")
        self.lbl_meta.pack(fill="x", pady=6)

        right = ttk.Frame(mid); right.pack(side="right", fill="y", padx=8)

        # ===== 범용 5문항 =====
        ttk.Label(right, text="Q1) 제목").pack(anchor="w")
        self.ent_title = ttk.Entry(right, width=48)
        self.ent_title.pack(fill="x", pady=4)

        ttk.Label(right, text="Q2) 고객사 (미지정 시 '범용')").pack(anchor="w")
        self.ent_customer = ttk.Entry(right, width=48)
        self.ent_customer.pack(fill="x", pady=4)

        ttk.Label(right, text="Q3) 분류").pack(anchor="w")
        self.cmb_category = ttk.Combobox(right, values=CATEGORIES, state="readonly", width=45)
        self.cmb_category.pack(fill="x", pady=4)

        ttk.Label(right, text="Q4) 사진 설명 (핵심 의미/상황)").pack(anchor="w")
        self.txt_description = tk.Text(right, height=10, width=46, wrap="word")
        self.txt_description.pack(fill="x", pady=6)

        ttk.Label(right, text="Q5) 태그(쉼표로 구분)").pack(anchor="w")
        self.ent_tags = ttk.Entry(right, width=48)
        self.ent_tags.pack(fill="x", pady=4)

        # OCR 정보 표시 영역
        sep_ocr = ttk.Separator(right, orient="horizontal")
        sep_ocr.pack(fill="x", pady=(10, 4))
        
        ttk.Label(right, text="OCR 분석 결과", font=("Malgun Gothic", 9, "bold")).pack(anchor="w", pady=(4, 2))
        self.lbl_ocr_summary = ttk.Label(right, text="OCR 정보 없음", foreground="gray", font=("Malgun Gothic", 8))
        self.lbl_ocr_summary.pack(anchor="w", pady=2)
        
        self.txt_ocr = tk.Text(right, height=6, width=46, wrap="word", font=("Malgun Gothic", 8), state="disabled")
        self.txt_ocr.pack(fill="x", pady=(4, 0))

        self.lbl_status = ttk.Label(right, text="")
        self.lbl_status.pack(anchor="w", pady=10)

        self.lbl_saved = ttk.Label(right, text="", foreground="#2b6")
        self.lbl_saved.pack(anchor="w", pady=4)

        self.lbl_remaining = ttk.Label(right, text="", foreground="#444")
        self.lbl_remaining.pack(anchor="w", pady=6)

    def reload(self):
        self.records = self._load_filtered_queue()
        self.index = 0
        self._refresh()

    def prev(self):
        if not self.records:
            return
        self.index = max(0, self.index - 1)
        self._refresh()

    def next(self):
        if not self.records:
            return
        self.index = min(len(self.records) - 1, self.index + 1)
        self._refresh()

    def goto_last(self):
        if not self.records:
            return
        self.index = len(self.records) - 1
        self._refresh()

    def finish(self):
        remain = self._count_remaining()
        if remain > 0:
            if not messagebox.askyesno("완료하기", f"아직 미처리 이미지가 {remain}개 남아있습니다.\n그래도 완료하고 창을 닫을까요?"):
                return
        else:
            messagebox.showinfo("완료", "모든 이미지가 처리된 것으로 보입니다. (저장기록 기준)")
        self.destroy()

    def _goto_next_after_action(self):
        if not self.records:
            return
        if self.index < len(self.records) - 1:
            self.index += 1
        self._refresh()

    def _refresh(self):
        n = len(self.records)
        if n == 0:
            self.lbl_pos.config(text="0/0 (큐가 비어있음)")
            self.canvas.delete("all")
            self.lbl_meta.config(text="images_queue.jsonl이 없거나 비어있습니다.")
            self._clear_form()
            self.lbl_status.config(text="")
            self.lbl_saved.config(text="")
            self.lbl_remaining.config(text="")
            return

        rec = self.records[self.index]
        ensure_dict(rec, "status", {"saved_label": False, "merged": False, "ignored": False})

        self.lbl_pos.config(text=f"{self.index+1}/{n}")

        img_path = rec["image_path"]
        try:
            img = Image.open(img_path)
            cw = int(self.canvas.winfo_width() or 680)
            ch = int(self.canvas.winfo_height() or 600)
            img.thumbnail((cw - 20, ch - 20))
            self._tk_img = ImageTk.PhotoImage(img)
            self.canvas.delete("all")
            self.canvas.create_image(cw // 2, ch // 2, image=self._tk_img)
        except Exception as e:
            self.canvas.delete("all")
            self.canvas.create_text(20, 20, anchor="nw", fill="white", text=f"이미지 로드 실패:\n{e}\n{img_path}")

        loc = rec.get("location") or {}
        loc_str = ""
        if "slide_index" in loc and loc.get("slide_index") is not None:
            loc_str = f"슬라이드 {int(loc['slide_index']) + 1}"
        elif loc:
            loc_str = str(loc)

        meta = (
            f"mode: {rec.get('mode')}\n"
            f"origin: {rec.get('origin')}\n"
            f"parent_doc_id: {rec.get('parent_doc_id')}\n"
            f"parent_source_path: {rec.get('parent_source_path')}\n"
            f"location: {loc_str if loc_str else rec.get('location')}\n"
            f"sha256: {rec.get('sha256')}\n"
            f"ignored: {rec.get('ignored')}\n"
        )
        self.lbl_meta.config(text=meta)
        
        # OCR 정보 표시 (있는 경우)
        ocr_info = rec.get("ocr")
        if ocr_info:
            ocr_text = ocr_info.get("text", "")
            ocr_conf = ocr_info.get("confidence", 0.0)
            ocr_words = ocr_info.get("word_count", 0)
            ocr_lines = ocr_info.get("line_count", 0)
            ocr_lang = ocr_info.get("language", "unknown")
            ocr_processed = ocr_info.get("processed_at", "")
            
            # OCR 요약 정보
            summary_parts = []
            if ocr_conf > 0:
                summary_parts.append(f"신뢰도: {ocr_conf:.1%}")
            if ocr_words > 0:
                summary_parts.append(f"단어: {ocr_words}")
            if ocr_lines > 0:
                summary_parts.append(f"라인: {ocr_lines}")
            if ocr_lang and ocr_lang != "unknown":
                summary_parts.append(f"언어: {ocr_lang}")
            
            self.lbl_ocr_summary.config(
                text=" | ".join(summary_parts) if summary_parts else "OCR 완료",
                foreground="blue"
            )
            
            # OCR 전체 텍스트
            self.txt_ocr.config(state="normal")
            self.txt_ocr.delete("1.0", "end")
            if ocr_text:
                self.txt_ocr.insert("1.0", ocr_text)
            else:
                self.txt_ocr.insert("1.0", "(OCR 텍스트 없음)")
            self.txt_ocr.config(state="disabled")
        else:
            self.lbl_ocr_summary.config(text="OCR 정보 없음", foreground="gray")
            self.txt_ocr.config(state="normal")
            self.txt_ocr.delete("1.0", "end")
            self.txt_ocr.insert("1.0", "(이미지에 OCR 분석 결과가 없습니다)")
            self.txt_ocr.config(state="disabled")

        # 폼 채우기
        lbl = rec.get("label", {})
        self._set_form(
            title=lbl.get("title", Path(img_path).stem),
            customer=lbl.get("customer", DEFAULT_CUSTOMER),
            category=lbl.get("category", CATEGORIES[0]),
            description=lbl.get("description", ""),
            tags=",".join(lbl.get("tags", []) or []),
        )

        st = rec.get("status", {})
        self.lbl_status.config(text=f"큐상태: labeled={rec.get('labeled')} / ignored={rec.get('ignored')}")
        self.lbl_saved.config(text=f"저장상태: 라벨저장={st.get('saved_label')} / 병합={st.get('merged')} / 무시={st.get('ignored')}")
        self.lbl_remaining.config(text=f"미처리(추정) 이미지: {self._count_remaining()}개")

    def _clear_form(self):
        self.ent_title.delete(0, "end")
        self.ent_customer.delete(0, "end")
        self.cmb_category.set(CATEGORIES[0])
        self.txt_description.delete("1.0", "end")
        self.ent_tags.delete(0, "end")
        # OCR 정보 초기화
        self.lbl_ocr_summary.config(text="OCR 정보 없음", foreground="gray")
        self.txt_ocr.config(state="normal")
        self.txt_ocr.delete("1.0", "end")
        self.txt_ocr.insert("1.0", "(OCR 정보 없음)")
        self.txt_ocr.config(state="disabled")

    def _set_form(self, title: str, customer: str, category: str, description: str, tags: str):
        self.ent_title.delete(0, "end")
        self.ent_title.insert(0, title or "")

        self.ent_customer.delete(0, "end")
        self.ent_customer.insert(0, customer or DEFAULT_CUSTOMER)

        if category in CATEGORIES:
            self.cmb_category.set(category)
        else:
            self.cmb_category.set(CATEGORIES[0])

        self.txt_description.delete("1.0", "end")
        self.txt_description.insert("end", description or "")

        self.ent_tags.delete(0, "end")
        self.ent_tags.insert(0, tags or "")

    def _update_rec_from_form(self, rec: dict) -> dict:
        title = (self.ent_title.get() or "").strip()
        customer = normalize_customer(self.ent_customer.get())
        category = self.cmb_category.get() if self.cmb_category.get() else CATEGORIES[0]
        description = (self.txt_description.get("1.0", "end") or "").strip()
        tags = normalize_tags(self.ent_tags.get())

        rec["label"] = {
            "title": title if title else Path(rec.get("image_path","")).stem,
            "customer": customer,
            "category": category,
            "description": description,
            "tags": tags,
        }
        return rec

    def _append_label_history(self, rec: dict):
        out = {
            "image_id": rec["image_id"],
            "image_path": rec["image_path"],
            "origin": rec.get("origin"),
            "mode": rec.get("mode"),
            "parent_doc_id": rec.get("parent_doc_id"),
            "parent_source_path": rec.get("parent_source_path"),
            "location": rec.get("location", {}),
            "label": rec.get("label", {}),
            "ignored": rec.get("ignored", False),
            "saved_at": now_iso(),
        }
        append_jsonl(self.labels_path, out)

    def _merge_metadata(self, rec: dict) -> bool:
        label = rec.get("label", {}) or {}
        image_item = {
            "image_id": rec["image_id"],
            "image_path": rec["image_path"],
            "origin": rec.get("origin", ""),
            "mode": rec.get("mode", "standalone"),
            "location": rec.get("location", {}),
            "title": label.get("title", ""),
            "customer": normalize_customer(label.get("customer")),
            "category": label.get("category", CATEGORIES[0]),
            "description": label.get("description", ""),
            "tags": label.get("tags", []),
            "ignored": rec.get("ignored", False),
            "updated_at": now_iso(),
        }
        
        # OCR 정보 추가 (있는 경우)
        if "ocr" in rec:
            image_item["ocr"] = rec["ocr"]

        if rec.get("mode") == "doc_embed" and rec.get("parent_doc_id"):
            upsert_doc_images_metadata(self.out_root, rec["parent_doc_id"], image_item)
            return True

        # standalone
        upsert_standalone_images_metadata(self.out_root, image_item)
        return True

    def save_current(self):
        if not self.records:
            return

        rec = self.records[self.index]
        rec = self._update_rec_from_form(rec)

        self._append_label_history(rec)
        merged = self._merge_metadata(rec)

        ensure_dict(rec, "status", {"saved_label": False, "merged": False, "ignored": False})
        rec["status"]["saved_label"] = True
        rec["status"]["merged"] = bool(merged)
        rec["status"]["ignored"] = bool(rec.get("ignored", False))
        rec["labeled"] = True  # 저장=처리로 간주

        self.records[self.index] = rec
        self._persist_queue()
        self._goto_next_after_action()

    def mark_ignored_and_next(self):
        if not self.records:
            return

        rec = self.records[self.index]
        rec = self._update_rec_from_form(rec)

        rec["ignored"] = True
        rec["labeled"] = True

        self._append_label_history(rec)
        merged = self._merge_metadata(rec)

        ensure_dict(rec, "status", {"saved_label": False, "merged": False, "ignored": False})
        rec["status"]["saved_label"] = True
        rec["status"]["merged"] = bool(merged)
        rec["status"]["ignored"] = True

        self.records[self.index] = rec
        self._persist_queue()
        self._goto_next_after_action()


class DocumentTopicWindow(tk.Toplevel):
    """문서 주제문장 입력 창"""
    def __init__(self, master, doc_name: str):
        super().__init__(master)
        self.title("문서 주제문장 입력")
        self.geometry("600x280")
        self.confirmed = False
        self.doc_name = doc_name
        
        self.topic_sentence = tk.StringVar()
        
        self._build_ui()
        self.transient(master)
        self.grab_set()
        self.focus_set()
        
    def _build_ui(self):
        pad = {"padx": 15, "pady": 10}
        
        # 제목
        title_frame = ttk.Frame(self)
        title_frame.pack(fill="x", pady=(15, 10))
        ttk.Label(title_frame, text="문서 주제문장 입력 (선택사항)", font=("Malgun Gothic", 11, "bold")).pack()
        
        # 문서명 표시
        doc_frame = ttk.Frame(self)
        doc_frame.pack(fill="x", padx=20, pady=(0, 10))
        ttk.Label(doc_frame, text=f"문서: {self.doc_name}", font=("Malgun Gothic", 9), foreground="gray").pack(anchor="w")
        
        # 설명
        desc_frame = ttk.Frame(self)
        desc_frame.pack(fill="x", padx=20, pady=(0, 10))
        desc_text = "주제문장을 입력하면 태그 우선순위와 장르 분류에 활용됩니다.\n입력하지 않으면 기존 오토태깅이 적용됩니다."
        ttk.Label(desc_frame, text=desc_text, foreground="gray", justify="left", font=("Malgun Gothic", 9)).pack(anchor="w")
        
        # 입력 필드
        input_frame = ttk.LabelFrame(self, text="주제문장", padding=10)
        input_frame.pack(fill="both", expand=True, padx=20, pady=10)
        
        self.text_entry = tk.Text(input_frame, height=4, wrap="word", font=("Malgun Gothic", 10))
        self.text_entry.pack(fill="both", expand=True)
        
        # 버튼
        btn_frame = ttk.Frame(self)
        btn_frame.pack(fill="x", padx=20, pady=(10, 15))
        ttk.Button(btn_frame, text="건너뛰기 (오토태깅 사용)", command=self.skip).pack(side="left", padx=5)
        ttk.Button(btn_frame, text="확인", command=self.confirm).pack(side="right", padx=5)
        
    def confirm(self):
        self.topic_sentence.set(self.text_entry.get("1.0", "end-1c").strip())
        self.confirmed = True
        self.destroy()
        
    def skip(self):
        self.topic_sentence.set("")
        self.confirmed = True
        self.destroy()


class ChunkingSettingsWindow(tk.Toplevel):
    """청킹 설정 창"""
    def __init__(self, master):
        super().__init__(master)
        self.title("청킹 설정")
        self.geometry("480x320")
        self.confirmed = False
        
        self.target_chars = tk.IntVar(value=1000)
        self.min_chars = tk.IntVar(value=250)
        self.overlap_chars = tk.IntVar(value=120)
        
        self._build_ui()
        self.transient(master)
        self.grab_set()
        
    def _build_ui(self):
        pad = {"padx": 15, "pady": 10}
        
        # 제목
        title_frame = ttk.Frame(self)
        title_frame.pack(fill="x", pady=(15, 10))
        ttk.Label(title_frame, text="청킹 파라미터 설정", font=("Malgun Gothic", 11, "bold")).pack()
        
        # 설정 항목들
        main_frame = ttk.Frame(self)
        main_frame.pack(fill="both", expand=True, padx=20, pady=10)
        
        # target_chars
        row1 = ttk.Frame(main_frame)
        row1.pack(fill="x", pady=8)
        ttk.Label(row1, text="목표 청크 크기 (target_chars):", width=22, anchor="w").pack(side="left")
        ttk.Spinbox(row1, from_=200, to=5000, textvariable=self.target_chars, width=12).pack(side="left", padx=8)
        ttk.Label(row1, text="자", foreground="gray").pack(side="left")
        
        # min_chars
        row2 = ttk.Frame(main_frame)
        row2.pack(fill="x", pady=8)
        ttk.Label(row2, text="최소 청크 크기 (min_chars):", width=22, anchor="w").pack(side="left")
        ttk.Spinbox(row2, from_=50, to=2000, textvariable=self.min_chars, width=12).pack(side="left", padx=8)
        ttk.Label(row2, text="자", foreground="gray").pack(side="left")
        
        # overlap_chars
        row3 = ttk.Frame(main_frame)
        row3.pack(fill="x", pady=8)
        ttk.Label(row3, text="청크 간 겹침 (overlap_chars):", width=22, anchor="w").pack(side="left")
        ttk.Spinbox(row3, from_=0, to=500, textvariable=self.overlap_chars, width=12).pack(side="left", padx=8)
        ttk.Label(row3, text="자", foreground="gray").pack(side="left")
        
        # 설명
        desc_frame = ttk.Frame(main_frame)
        desc_frame.pack(fill="x", pady=(15, 10))
        desc_text = "• 목표 크기: 각 청크의 목표 문자 수\n• 최소 크기: 청크의 최소 문자 수\n• 겹침: 인접 청크 간 겹치는 문자 수"
        ttk.Label(desc_frame, text=desc_text, foreground="gray", justify="left", font=("Malgun Gothic", 9)).pack(anchor="w")
        
        # 버튼
        btn_frame = ttk.Frame(self)
        btn_frame.pack(fill="x", padx=20, pady=(10, 15))
        ttk.Button(btn_frame, text="취소", command=self.cancel).pack(side="right", padx=5)
        ttk.Button(btn_frame, text="확인", command=self.confirm).pack(side="right", padx=5)
        
    def confirm(self):
        self.confirmed = True
        self.destroy()
        
    def cancel(self):
        self.confirmed = False
        self.destroy()


class DocumentSelectWindow(tk.Toplevel):
    def __init__(self, master, out_root):
        super().__init__(master)
        self.title("문서 선택")
        self.geometry("760x460")
        self.out_root = out_root
        self.selected = None

        self.queue = []
        self.payload_docs = []   # workspace_payload.jsonl 로딩 결과
        self.docs = {}           # doc_id -> {name,total,unlabeled,missing}
        self.keys = []

        self._load_from_disk()
        self._build_ui()

    # -------------------------
    # Load
    # -------------------------
    def _load_from_disk(self):
        out = Path(self.out_root)

        # 1) 이미지 큐(카운트용)
        self.queue = read_jsonl(str(out / "images_queue.jsonl"))

        # 2) 문서 목록(진짜 “문서함”의 근거)
        self.payload_docs = read_jsonl(str(out / "workspace_payload.jsonl"))

        self.docs = self._build_doc_index()

    def _build_doc_index(self):
        """
        문서함 목록은 workspace_payload.jsonl 기준으로 구성한다.
        - images_queue.jsonl은 문서별 이미지(미라벨/전체) 카운트 산출용으로만 사용
        - 추출 이미지가 0인 문서도 목록에 반드시 표시됨
        """
        index = {}

        # A) workspace_payload 기반 문서 목록 생성
        for r in self.payload_docs:
            doc_id = r.get("doc_id")
            src = r.get("source_path") or ""
            if not doc_id:
                continue

            name = Path(src).name if src else (r.get("title") or f"[doc_id:{doc_id}]")
            # 원본 문서 파일 존재 여부(참고)
            src_exists = bool(src) and Path(src).exists()
            if not src_exists:
                name = f"{name}  (원본 없음)"

            index.setdefault(doc_id, {
                "name": name,
                "total": 0,
                "unlabeled": 0,
                "missing": 0,   # 이미지 파일 없음 카운트
            })

        # B) images_queue로 문서별 이미지 카운트 합산
        #    (문서에 이미지가 0이어도 index에는 이미 들어가 있음)
        for img in self.queue:
            mode = img.get("mode")
            if mode == "standalone":
                key = "__STANDALONE__"
                index.setdefault(key, {
                    "name": "[순수 이미지] standalone",
                    "total": 0,
                    "unlabeled": 0,
                    "missing": 0,
                })
            else:
                key = img.get("parent_doc_id") or "(unknown_doc)"
                index.setdefault(key, {
                    "name": f"[doc_id:{key}]",
                    "total": 0,
                    "unlabeled": 0,
                    "missing": 0,
                })

            img_path = img.get("image_path")
            img_exists = bool(img_path) and Path(img_path).exists()

            index[key]["total"] += 1
            if not img.get("labeled"):
                index[key]["unlabeled"] += 1
            if not img_exists:
                index[key]["missing"] += 1

        # workspace_payload가 비어있고, standalone도 없고, 큐도 없는 “진짜 0”일 때 안내용
        return index

    # -------------------------
    # UI
    # -------------------------
    def _build_ui(self):
        # 제목
        title_frame = ttk.Frame(self)
        title_frame.pack(fill="x", pady=(15, 10))
        ttk.Label(title_frame, text="라벨링할 문서를 선택하세요", font=("Malgun Gothic", 11, "bold")).pack()

        # 리스트박스 프레임
        list_frame = ttk.LabelFrame(self, text="문서 목록", padding=10)
        list_frame.pack(fill="both", expand=True, padx=12, pady=8)
        self.listbox = tk.Listbox(list_frame, height=16, font=("Malgun Gothic", 9))
        self.listbox.pack(fill="both", expand=True)

        # 버튼 행
        btn_row = ttk.Frame(self)
        btn_row.pack(fill="x", padx=12, pady=(8, 12))

        ttk.Button(btn_row, text="문서함 갱신", command=self.refresh, width=16).pack(side="left", padx=4)
        ttk.Button(btn_row, text="입력폴더 재동기화", command=self.resync_with_input, width=20).pack(side="left", padx=4)
        ttk.Button(btn_row, text="선택", command=self.confirm, width=12).pack(side="right", padx=4)

        self._render_list()

    def _render_list(self):
        self.listbox.delete(0, "end")
        self.keys = []

        items = list(self.docs.items())
        items.sort(key=lambda kv: (kv[1].get("name", ""), str(kv[0])))

        for k, v in items:
            label = f"{v['name']}  (미라벨 {v['unlabeled']} / 전체 {v['total']} / 파일없음 {v.get('missing',0)})"
            self.listbox.insert("end", label)
            self.keys.append(k)

        if not self.keys:
            self.listbox.insert("end", "(표시할 문서/이미지가 없습니다)")

    def _prune_missing_against_input(self, also_prune_payload=False, also_delete_doc_dirs=False):
        """
        입력 폴더에 없는 문서의 관련 데이터를 정리합니다.
        
        Args:
            also_prune_payload: workspace_payload.jsonl에서도 제거할지
            also_delete_doc_dirs: output/<doc_id>/ 디렉터리도 삭제할지
        """
        app = self.master
        input_dir = (app.input_dir.get() or "").strip()
        out_dir = self.out_root
        
        if not input_dir or not Path(input_dir).exists():
            return  # 입력 폴더가 없으면 정리할 것도 없음
        
        # 1) 입력 폴더의 현재 문서 ID 수집
        docs, _ = scan_input(input_dir)
        current_doc_ids = {doc_id_for(p) for p in docs}
        
        # 2) 이미지 큐 정리
        queue_path = Path(out_dir) / "images_queue.jsonl"
        if queue_path.exists():
            queue = read_jsonl(str(queue_path))
            pruned_queue = []
            removed_queue = 0
            
            for r in queue:
                if r.get("mode") == "standalone":
                    pruned_queue.append(r)  # 순수 이미지는 유지
                    continue
                parent_doc_id = r.get("parent_doc_id")
                if parent_doc_id and parent_doc_id in current_doc_ids:
                    pruned_queue.append(r)  # 입력 폴더에 있는 문서의 이미지는 유지
                else:
                    removed_queue += 1  # 입력 폴더에 없는 문서의 이미지는 제거
            
            write_jsonl(str(queue_path), pruned_queue)
        
        # 3) workspace_payload.jsonl 정리 (옵션)
        if also_prune_payload:
            payload_path = Path(out_dir) / "workspace_payload.jsonl"
            if payload_path.exists():
                payload = read_jsonl(str(payload_path))
                pruned_payload = [r for r in payload if r.get("doc_id") in current_doc_ids]
                write_jsonl(str(payload_path), pruned_payload)
        
        # 4) 디렉터리 삭제 (옵션)
        if also_delete_doc_dirs:
            out_path = Path(out_dir)
            if out_path.exists():
                for doc_id_dir in out_path.iterdir():
                    if doc_id_dir.is_dir() and doc_id_dir.name not in ("core", "metadata_docs"):
                        if doc_id_dir.name not in current_doc_ids:
                            try:
                                shutil.rmtree(doc_id_dir)
                            except Exception:
                                pass  # 삭제 실패해도 계속 진행

    def refresh(self):
        try:
            self._prune_missing_against_input(
                also_prune_payload=True,
                also_delete_doc_dirs=True
            )
        except Exception:
            pass

        self._load_from_disk()
        self._render_list()

    def confirm(self):
        sel = self.listbox.curselection()
        if not sel or not self.keys:
            messagebox.showwarning("선택 필요", "문서를 선택하세요.")
            return
        self.selected = self.keys[sel[0]]
        self.destroy()

    # -------------------------
    # Resync (네 기존 코드 유지하되, 끝에 refresh만 호출하면 됨)
    # -------------------------
    def resync_with_input(self):
        """
        입력 폴더 기준으로 재동기화:
        - (기존 네 코드 그대로 두되)
        - 마지막에 self.refresh()만 호출하면 문서함이 payload 기반으로 확실히 갱신됨
        """
        try:
            app = self.master

            input_dir = (app.input_dir.get() or "").strip()
            out_dir = (app.output_dir.get() or "").strip()
            tmp_dir = (app.tmp_dir.get() or "").strip()
            soffice = (app.soffice_path.get() or "").strip()
            scope = (app.inference_scope.get() or "").strip()

            if not input_dir or not Path(input_dir).exists():
                messagebox.showerror("오류", "입력 폴더가 설정되어 있지 않거나 존재하지 않습니다.")
                return
            if not out_dir:
                messagebox.showerror("오류", "출력 폴더가 유효하지 않습니다.")
                return

            Path(out_dir).mkdir(parents=True, exist_ok=True)
            Path(tmp_dir).mkdir(parents=True, exist_ok=True)

            docs, _ = scan_input(input_dir)
            current_doc_ids = {doc_id_for(p) for p in docs}

            # (1) 큐 프루닝
            queue_path = Path(out_dir) / "images_queue.jsonl"
            queue = read_jsonl(str(queue_path))

            pruned = []
            removed = 0
            for r in queue:
                if r.get("mode") == "standalone":
                    pruned.append(r)
                    continue
                if r.get("parent_doc_id") in current_doc_ids:
                    pruned.append(r)
                else:
                    removed += 1
            write_jsonl(str(queue_path), pruned)

            # (2) 새 문서 인제스트
            payload_path = Path(out_dir) / "workspace_payload.jsonl"
            existing_doc_ids = load_existing_doc_ids(str(payload_path))
            existing_image_ids = load_existing_image_ids(str(queue_path))

            ingested = 0
            queued_imgs = 0
            failed = 0

            for f in docs:
                did = doc_id_for(f)
                if did in existing_doc_ids:
                    continue
                try:
                    if Path(f).suffix.lower() == ".ppt" and (not soffice or not Path(soffice).exists()):
                        raise FileNotFoundError(f"soffice.exe 경로가 유효하지 않습니다: {soffice}")

                    # 메모리 최적화 설정 전달
                    use_semantic = getattr(self, 'use_semantic_adjustment', tk.BooleanVar(value=False)).get()
                    enable_ocr = getattr(self, 'enable_ocr', tk.BooleanVar(value=True)).get()
                    max_ocr = getattr(self, 'max_ocr_images', tk.IntVar(value=5)).get()
                    
                    res = process_one_doc(f, out_dir, tmp_dir, soffice, scope, 
                                         use_semantic=use_semantic, 
                                         enable_ocr=enable_ocr, 
                                         max_ocr_images=max_ocr)
                    ingested += 1
                    existing_doc_ids.add(res["doc_id"])

                    for img_rec in res.get("extracted_images", []):
                        loc = {}
                        if img_rec.get("slide_index") is not None:
                            loc = {"slide_index": img_rec.get("slide_index")}

                        rec = register_image_for_review(
                            out_root=out_dir,
                            image_path=img_rec["image_path"],
                            origin=img_rec.get("origin", "doc_embed"),
                            mode="doc_embed",
                            parent_doc_id=res["doc_id"],
                            parent_path=str(Path(f).resolve()),
                            location=loc,
                            existing_image_ids=existing_image_ids,
                            ocr_info=img_rec.get("ocr"),  # OCR 정보 전달
                        )
                        if not rec.get("skipped"):
                            queued_imgs += 1

                except Exception:
                    failed += 1

            # ✅ 여기서 payload 기반 문서함 갱신
            self.refresh()

            messagebox.showinfo(
                "재동기화 완료",
                f"입력폴더 재동기화 완료\n"
                f"- 제거된 큐 항목: {removed}개\n"
                f"- 새 문서 인제스트: {ingested}개\n"
                f"- 새 이미지 큐잉: {queued_imgs}개\n"
                f"- 인제스트 실패: {failed}개"
            )

        except Exception as e:
            messagebox.showerror("오류", f"재동기화 중 오류: {e}")

    


# =========================
# Main App
# =========================
class App(tk.Tk):
    def __init__(self):
        try:
            super().__init__()
            self.title("Metadata Prep - Docs + Standalone Images (범용 5문항)")
            self.geometry("1040x820")

            self.input_dir = tk.StringVar()
            self.output_dir = tk.StringVar(value=str(Path.cwd() / "output"))
            self.tmp_dir = tk.StringVar(value=str(Path.cwd() / "_tmp_convert"))
            self.soffice_path = tk.StringVar(value=r"C:\Program Files\LibreOffice\program\soffice.exe")
            self.inference_scope = tk.StringVar(value="internal")

            self.stop_flag = threading.Event()
            self.worker = None
            self._show_topic_input = False  # 주제문장 입력 모드 (기본값: False)
            
            # 메모리 최적화 설정
            self.use_semantic_adjustment = tk.BooleanVar(value=False)  # 의미 보정 사용 여부 (기본값: False로 메모리 절약)
            self.enable_ocr = tk.BooleanVar(value=True)  # OCR 사용 여부
            self.max_ocr_images = tk.IntVar(value=5)  # 최대 OCR 처리 이미지 수 (기본값: 5개로 제한)

            self._build_ui()
        except Exception as e:
            import traceback
            print("=" * 60)
            print("App 초기화 중 오류 발생:")
            print("=" * 60)
            traceback.print_exc()
            print("=" * 60)
            raise
    
    def _toggle_topic_input_mode(self):
        """주제문장 입력 모드 토글"""
        self._show_topic_input = self.topic_input_mode.get()

    def _build_ui(self):
        pad = {"padx": 12, "pady": 8}

        # 메인 프레임
        frm = ttk.Frame(self)
        frm.pack(fill="both", expand=True, padx=10, pady=10)

        # 설정 섹션
        settings_frame = ttk.LabelFrame(frm, text="설정", padding=10)
        settings_frame.pack(fill="x", pady=(0, 10))

        row = ttk.Frame(settings_frame); row.pack(fill="x", pady=4)
        ttk.Label(row, text="입력 폴더(문서/이미지 혼합)", width=22, anchor="w").pack(side="left")
        ttk.Entry(row, textvariable=self.input_dir, width=70, font=("Malgun Gothic", 9)).pack(side="left", padx=8)
        ttk.Button(row, text="찾기", command=self.pick_input, width=8).pack(side="left")

        row = ttk.Frame(settings_frame); row.pack(fill="x", pady=4)
        ttk.Label(row, text="출력 폴더", width=22, anchor="w").pack(side="left")
        ttk.Entry(row, textvariable=self.output_dir, width=70, font=("Malgun Gothic", 9)).pack(side="left", padx=8)
        ttk.Button(row, text="찾기", command=self.pick_output, width=8).pack(side="left")

        row = ttk.Frame(settings_frame); row.pack(fill="x", pady=4)
        ttk.Label(row, text="변환 임시폴더", width=22, anchor="w").pack(side="left")
        ttk.Entry(row, textvariable=self.tmp_dir, width=70, font=("Malgun Gothic", 9)).pack(side="left", padx=8)
        ttk.Button(row, text="찾기", command=self.pick_tmp, width=8).pack(side="left")

        row = ttk.Frame(settings_frame); row.pack(fill="x", pady=4)
        ttk.Label(row, text="soffice.exe (ppt 변환)", width=22, anchor="w").pack(side="left")
        ttk.Entry(row, textvariable=self.soffice_path, width=62, font=("Malgun Gothic", 9)).pack(side="left", padx=8)
        ttk.Button(row, text="파일 선택", command=self.pick_soffice, width=10).pack(side="left")

        row = ttk.Frame(settings_frame); row.pack(fill="x", pady=4)
        ttk.Label(row, text="inference_scope", width=22, anchor="w").pack(side="left")
        ttk.Combobox(
            row,
            textvariable=self.inference_scope,
            values=["public", "internal", "restricted"],
            width=18,
            state="readonly",
            font=("Malgun Gothic", 9)
        ).pack(side="left", padx=8)
        
        # 주제문장 입력 모드 체크박스
        row = ttk.Frame(settings_frame); row.pack(fill="x", pady=4)
        self.topic_input_mode = tk.BooleanVar(value=False)
        ttk.Checkbutton(
            row,
            text="문서 처리 시 주제문장 입력 요청",
            variable=self.topic_input_mode,
            command=self._toggle_topic_input_mode
        ).pack(side="left")
        
        # 메모리 최적화 설정
        row = ttk.Frame(settings_frame); row.pack(fill="x", pady=4)
        ttk.Label(row, text="메모리 최적화:", width=22, anchor="w").pack(side="left")
        ttk.Checkbutton(
            row,
            text="의미 보정 사용",
            variable=self.use_semantic_adjustment
        ).pack(side="left", padx=5)
        ttk.Checkbutton(
            row,
            text="OCR 활성화",
            variable=self.enable_ocr
        ).pack(side="left", padx=5)
        ttk.Label(row, text="최대 OCR 이미지:", width=12, anchor="w").pack(side="left", padx=(20, 0))
        ocr_max_spin = ttk.Spinbox(row, from_=0, to=50, textvariable=self.max_ocr_images, width=8)
        ocr_max_spin.pack(side="left", padx=5)

        # 첫 번째 줄: 문서 처리 관련 버튼들
        row1 = ttk.Frame(frm); row1.pack(fill="x", **pad)
        self.btn_start = ttk.Button(row1, text="문서/폴더 일괄 처리 시작(누적)", command=self.start)
        self.btn_start.pack(side="left")
        self.btn_stop = ttk.Button(row1, text="중지", command=self.stop, state="disabled")
        self.btn_stop.pack(side="left", padx=8)
        self.btn_add_imgs = ttk.Button(row1, text="순수 이미지 추가(폴더/파일)", command=self.add_standalone_images)
        self.btn_add_imgs.pack(side="left", padx=12)
        self.btn_review = ttk.Button(row1, text="이미지 라벨링(범용 5문항)", command=self.open_review, state="disabled")
        self.btn_review.pack(side="left", padx=12)
        self.btn_record = ttk.Button(row1, text='업무 레코드/보고서 생성(수동 입력)', command=self.open_record_builder)
        self.btn_record.pack(side='left', padx=12)

        # 두 번째 줄: 청킹 및 검색 관련 버튼들
        row2 = ttk.Frame(frm); row2.pack(fill="x", **pad)
        self.btn_export = ttk.Button(row2, text="정제 메타DOCX 생성(전체)", command=self.export_all_metadata_docs)
        self.btn_export.pack(side="left", padx=12)
        self.btn_chunking = ttk.Button(row2, text="문서 청킹 처리(선택)", command=self.open_chunking)
        self.btn_chunking.pack(side="left", padx=12)
        self.btn_chunking_viewer = ttk.Button(row2, text="청킹 결과 보기", command=self.open_chunking_viewer)
        self.btn_chunking_viewer.pack(side="left", padx=12)
        self.btn_search = ttk.Button(row2, text="문서 검색", command=self.open_document_search)
        self.btn_search.pack(side="left", padx=12)
        self.btn_tagging_viewer = ttk.Button(row2, text="태깅 결과 보기", command=self.open_tagging_viewer)
        self.btn_tagging_viewer.pack(side="left", padx=12)
        self.btn_classification_dashboard = ttk.Button(row2, text="문서 분류 대시보드", command=self.open_classification_dashboard)
        self.btn_classification_dashboard.pack(side="left", padx=12)

        row = ttk.Frame(frm); row.pack(fill="x", **pad)
        self.prog = ttk.Progressbar(row, length=680, mode="determinate")
        self.prog.pack(side="left", padx=8)
        self.lbl_prog = ttk.Label(row, text="대기 중")
        self.lbl_prog.pack(side="left")

        # 통계 그래프 영역 (처리 완료 후 표시)
        self.stats_frame = ttk.LabelFrame(frm, text="문서 통계")
        # 초기에는 숨김 상태, 처리 완료 시 표시
        
        ttk.Label(frm, text="로그").pack(anchor="w", padx=8)
        self.txt = tk.Text(frm, height=28, wrap="word")
        self.txt.pack(fill="both", expand=True, padx=8, pady=6)

    def log(self, msg: str):
        self.txt.insert("end", msg + "\n")
        self.txt.see("end")

    def pick_input(self):
        d = filedialog.askdirectory(title="입력 폴더 선택")
        if d:
            self.input_dir.set(d)

    def pick_output(self):
        d = filedialog.askdirectory(title="출력 폴더 선택")
        if d:
            self.output_dir.set(d)

    def pick_tmp(self):
        d = filedialog.askdirectory(title="임시 폴더 선택")
        if d:
            self.tmp_dir.set(d)

    def pick_soffice(self):
        f = filedialog.askopenfilename(title="soffice.exe 선택", filetypes=[("exe", "*.exe")])
        if f:
            self.soffice_path.set(f)

    def _ui_progress(self, cur, total, label):
        self.lbl_prog.config(text=f"{label} ({cur}/{total})")
        self.update_idletasks()

    def start(self):
        in_dir = self.input_dir.get().strip()
        out_dir = self.output_dir.get().strip()
        tmp_dir = self.tmp_dir.get().strip()

        if not in_dir or not Path(in_dir).exists():
            messagebox.showerror("오류", "입력 폴더를 선택하세요.")
            return

        Path(out_dir).mkdir(parents=True, exist_ok=True)
        Path(tmp_dir).mkdir(parents=True, exist_ok=True)

        self.stop_flag.clear()
        self.btn_start.config(state="disabled")
        self.btn_stop.config(state="normal")
        self.btn_review.config(state="disabled")
        self.txt.delete("1.0", "end")
        self.log("스캔 시작... (누적 모드: 기존 데이터 유지 + 중복 스킵)")

        self.worker = threading.Thread(target=self._run_batch, daemon=True)
        self.worker.start()

    def stop(self):
        self.stop_flag.set()
        self.log("중지 요청됨... (현재 파일 처리 후 종료)")
        self.btn_stop.config(state="disabled")

    def _run_batch(self):
        try:
            in_dir = self.input_dir.get().strip()
            out_dir = self.output_dir.get().strip()
            tmp_dir = self.tmp_dir.get().strip()
            soffice = self.soffice_path.get().strip()
            scope = self.inference_scope.get().strip()

            docs, imgs = scan_input(in_dir)

            payload_path = str(Path(out_dir) / "workspace_payload.jsonl")
            queue_path = str(Path(out_dir) / "images_queue.jsonl")
            existing_doc_ids = load_existing_doc_ids(payload_path)
            existing_image_ids = load_existing_image_ids(queue_path)

            total = len(docs) + len(imgs)
            self.prog["maximum"] = max(1, total)
            self.prog["value"] = 0
            self._ui_progress(0, total, "스캔 완료")

            self.log(f"문서 {len(docs)}개, 이미지파일 {len(imgs)}개 발견")

            ok_docs = 0
            fail_docs = 0

            for i, f in enumerate(docs, start=1):
                if self.stop_flag.is_set():
                    self.log("중지 플래그 감지 -> 종료")
                    break

                self._ui_progress(self.prog["value"], total, f"문서 처리: {Path(f).name}")
                self.log(f"[DOC {i}/{len(docs)}] {f}")

                try:
                    if Path(f).suffix.lower() == ".ppt" and not Path(soffice).exists():
                        raise FileNotFoundError(f"soffice.exe 경로가 유효하지 않습니다: {soffice}")

                    did = doc_id_for(f)
                    if did in existing_doc_ids:
                        self.log(f"  ↩️ SKIP (already ingested) doc_id={did}")
                        self.prog["value"] += 1
                        continue

                    # 주제문장 입력 UI 표시 (선택적)
                    topic_sentence = None
                    if hasattr(self, '_show_topic_input') and self._show_topic_input:
                        topic_window = DocumentTopicWindow(self, Path(f).name)
                        self.wait_window(topic_window)
                        if topic_window.confirmed:
                            topic_sentence = topic_window.topic_sentence.get().strip() or None
                    
                    # 메모리 최적화 설정 전달
                    use_semantic = getattr(self, 'use_semantic_adjustment', tk.BooleanVar(value=False)).get()
                    enable_ocr = getattr(self, 'enable_ocr', tk.BooleanVar(value=True)).get()
                    max_ocr = getattr(self, 'max_ocr_images', tk.IntVar(value=5)).get()
                    
                    res = process_one_doc(f, out_dir, tmp_dir, soffice, scope, 
                                         topic_sentence=topic_sentence,
                                         use_semantic=use_semantic, 
                                         enable_ocr=enable_ocr, 
                                         max_ocr_images=max_ocr)
                    ok_docs += 1
                    self.log(f"  ✅ OK doc_id={res['doc_id']}")
                    if topic_sentence:
                        self.log(f"    📝 주제문장: {topic_sentence[:50]}...")
                    existing_doc_ids.add(res["doc_id"])

                    for img_rec in res.get("extracted_images", []):
                        loc = {}
                        if img_rec.get("slide_index") is not None:
                            loc = {"slide_index": img_rec.get("slide_index")}

                        rec = register_image_for_review(
                            out_root=out_dir,
                            image_path=img_rec["image_path"],
                            origin=img_rec.get("origin", "doc_embed"),
                            mode="doc_embed",
                            parent_doc_id=res["doc_id"],
                            parent_path=str(Path(f).resolve()),
                            location=loc,
                            existing_image_ids=existing_image_ids,
                            ocr_info=img_rec.get("ocr"),  # OCR 정보 전달
                        )
                        if rec.get("skipped"):
                            self.log(f"    ↩️ IMG SKIP (already queued): {rec['image_path']}")
                        else:
                            self.log(f"    🖼️ IMG queued: {Path(rec['image_path']).name}")

                except Exception as e:
                    fail_docs += 1
                    self.log(f"  ❌ FAIL: {e}")
                    append_jsonl(str(Path(out_dir) / "failed_docs.jsonl"), {
                        "path": f,
                        "error": str(e),
                        "trace": traceback.format_exc(),
                        "time": now_iso()
                    })

                self.prog["value"] += 1

            # 폴더에 “순수 이미지 파일”이 섞여있으면 기본은 standalone로 큐잉
            for j, img_path in enumerate(imgs, start=1):
                if self.stop_flag.is_set():
                    break
                self._ui_progress(self.prog["value"], total, f"이미지 등록(standalone): {Path(img_path).name}")
                try:
                    rec = register_image_for_review(
                        out_root=out_dir,
                        image_path=img_path,
                        origin="folder_image",
                        mode="standalone",
                        parent_doc_id=None,
                        parent_path=None,
                        location={},
                        existing_image_ids=existing_image_ids,
                    )
                    if rec.get("skipped"):
                        self.log(f"[IMG {j}/{len(imgs)}] ↩️ SKIP (already queued): {img_path}")
                    else:
                        self.log(f"[IMG {j}/{len(imgs)}] ✅ queued: {img_path}")
                except Exception as e:
                    self.log(f"[IMG {j}/{len(imgs)}] ❌ FAIL: {e}")
                self.prog["value"] += 1

            self._ui_progress(self.prog["value"], total, "완료")
            self.log("\n=== 완료 ===")
            self.log(f"문서: 성공 {ok_docs}, 실패 {fail_docs}")
            self.log(f"출력 폴더: {out_dir}")
            self.log(f"- workspace_payload.jsonl (누적)")
            self.log(f"- images_queue.jsonl (누적 / 라벨링 대상)")
            self.log(f"- images_labels.jsonl (라벨 저장 이력)")
            self.log(f"- output/<doc_id>/images_metadata.json (문서별 병합)")
            self.log(f"- output/standalone_images_metadata.json (순수 이미지 병합)")
            self.log(f"- output/metadata_docs/<원본문서명>.docx (정제 결과물)")
            self.log(f"- output/metadata_docs/_IMAGE_DICTIONARY.docx (순수 이미지 사전)")

            queued = len(read_jsonl(str(Path(out_dir) / "images_queue.jsonl")))
            self.log(f"이미지 라벨링 대상(큐): {queued}개")
            if queued > 0:
                self.btn_review.config(state="normal")
            
            # 문서 통계 그래프 표시
            if ok_docs > 0:
                self._show_document_statistics(out_dir)

        finally:
            self.btn_start.config(state="normal")
            self.btn_stop.config(state="disabled")

    def _show_document_statistics(self, out_dir: str):
        """문서 처리 완료 후 통계 그래프 표시"""
        try:
            # 기존 통계 그래프 제거
            if hasattr(self, 'stats_frame') and hasattr(self.stats_frame, 'winfo_exists'):
                try:
                    if self.stats_frame.winfo_exists():
                        self.stats_frame.destroy()
                except:
                    pass
            
            if not MATPLOTLIB_AVAILABLE:
                return
            
            # workspace_payload.jsonl에서 모든 문서 로드
            payload_path = Path(out_dir) / "workspace_payload.jsonl"
            if not payload_path.exists():
                return
            
            payloads = read_jsonl(str(payload_path))
            if not payloads:
                return
            
            # 장르별 통계 계산
            genres = Counter()
            languages = Counter()
            total_tags = 0
            avg_confidence = 0.0
            
            for payload in payloads:
                doc_id = payload.get("doc_id")
                if not doc_id:
                    continue
                
                auto_tags = load_auto_tags(out_dir, doc_id) or {}
                genre = auto_tags.get("genre", "unknown")
                language = auto_tags.get("language", "unknown")
                tags_topk = auto_tags.get("tags_topk", [])
                genre_conf = auto_tags.get("genre_confidence", 0.0)
                
                genres[genre] += 1
                languages[language] += 1
                total_tags += len(tags_topk)
                avg_confidence += genre_conf
            
            total_docs = len(payloads)
            if total_docs > 0:
                avg_confidence /= total_docs
            else:
                return
            
            # 장르 한글명 매핑
            genre_names_ko = {
                "issue": "이슈/문제",
                "resolution": "해결책",
                "procedure": "절차/가이드",
                "report": "보고서",
                "policy": "정책",
                "communication": "의사소통",
                "plan": "계획",
                "contract": "계약/제안",
                "reference": "참고자료",
                "application": "신청서",
                "form": "양식",
                "maintenance": "유지보수",
                "guide": "가이드",
                "record": "기록",
                "unknown": "미분류"
            }
            
            # 통계 프레임 생성 위치 찾기 (로그 라벨 위)
            # frm을 찾기 위해 winfo_children를 사용
            frm = None
            for widget in self.winfo_children():
                if isinstance(widget, ttk.Frame):
                    # 내부에 로그 라벨이 있는지 확인
                    for child in widget.winfo_children():
                        if isinstance(child, tk.Text) and child == self.txt:
                            frm = widget
                            break
                    if frm:
                        break
            
            if frm is None:
                return
            
            # 통계 프레임 생성
            self.stats_frame = ttk.LabelFrame(frm, text="문서 통계")
            # 로그 라벨 바로 위에 배치 (로그 라벨을 찾아서 그 위에 배치)
            log_label = None
            for child in frm.winfo_children():
                if isinstance(child, ttk.Label):
                    text = child.cget("text")
                    if text == "로그":
                        log_label = child
                        break
            
            if log_label:
                self.stats_frame.pack(fill="x", padx=8, pady=(10, 5), before=log_label)
            else:
                self.stats_frame.pack(fill="x", padx=8, pady=(10, 5))
            
            # 그래프 생성
            fig = Figure(figsize=(10, 4), dpi=80)
            
            # 1. 장르별 분포
            ax1 = fig.add_subplot(1, 2, 1)
            if genres:
                genre_names = [genre_names_ko.get(g, g) for g in genres.keys()]
                genre_counts = list(genres.values())
                colors = plt.cm.Set3(range(len(genre_names)))
                ax1.barh(genre_names, genre_counts, color=colors, alpha=0.7)
                ax1.set_xlabel('문서 수', fontproperties=font_prop if font_prop else None, fontsize=9)
                ax1.set_title('장르별 문서 분포', fontproperties=font_prop if font_prop else None, fontsize=10, fontweight='bold')
                ax1.grid(axis='x', alpha=0.3)
                # 값 표시
                for i, v in enumerate(genre_counts):
                    ax1.text(v + 0.1, i, str(v), va='center', fontproperties=font_prop if font_prop else None, fontsize=8)
            else:
                ax1.text(0.5, 0.5, '데이터 없음', ha='center', va='center', transform=ax1.transAxes, fontproperties=font_prop if font_prop else None)
                ax1.set_title('장르별 문서 분포', fontproperties=font_prop if font_prop else None, fontsize=10, fontweight='bold')
            
            # 2. 언어별 분포
            ax2 = fig.add_subplot(1, 2, 2)
            if languages:
                lang_names = list(languages.keys())
                lang_counts = list(languages.values())
                colors = plt.cm.Pastel1(range(len(lang_names)))
                ax2.barh(lang_names, lang_counts, color=colors, alpha=0.7)
                ax2.set_xlabel('문서 수', fontproperties=font_prop if font_prop else None, fontsize=9)
                ax2.set_title('언어별 문서 분포', fontproperties=font_prop if font_prop else None, fontsize=10, fontweight='bold')
                ax2.grid(axis='x', alpha=0.3)
                # 값 표시
                for i, v in enumerate(lang_counts):
                    ax2.text(v + 0.1, i, str(v), va='center', fontproperties=font_prop if font_prop else None, fontsize=8)
            else:
                ax2.text(0.5, 0.5, '데이터 없음', ha='center', va='center', transform=ax2.transAxes, fontproperties=font_prop if font_prop else None)
                ax2.set_title('언어별 문서 분포', fontproperties=font_prop if font_prop else None, fontsize=10, fontweight='bold')
            
            plt.tight_layout()
            
            # Tkinter에 그래프 표시
            canvas = FigureCanvasTkAgg(fig, self.stats_frame)
            canvas.draw()
            canvas.get_tk_widget().pack(fill="both", expand=True, padx=5, pady=5)
            
            # 통계 정보 텍스트 표시
            avg_tags = total_tags / total_docs if total_docs > 0 else 0
            stats_text = f"전체 문서: {total_docs}개 | 평균 태그 수: {avg_tags:.1f}개 | 평균 장르 신뢰도: {avg_confidence:.2f}"
            stats_label = ttk.Label(self.stats_frame, text=stats_text, font=("Malgun Gothic", 9))
            stats_label.pack(pady=(0, 5))
            
        except Exception as e:
            self.log(f"⚠️ 통계 그래프 생성 실패: {e}")
            import traceback
            self.log(traceback.format_exc())

    def add_standalone_images(self):
        out_dir = self.output_dir.get().strip()
        if not out_dir:
            messagebox.showerror("오류", "출력 폴더를 먼저 지정하세요.")
            return
        Path(out_dir).mkdir(parents=True, exist_ok=True)

        queue_path = str(Path(out_dir) / "images_queue.jsonl")
        existing_image_ids = load_existing_image_ids(queue_path)

        # 폴더 또는 파일을 선택 (둘 중 하나)
        if messagebox.askyesno("순수 이미지 추가", "폴더로 추가할까요?\n(아니오 선택 시 파일로 추가)"):
            d = filedialog.askdirectory(title="이미지 폴더 선택")
            if not d:
                return
            img_files = []
            for root, _, files in os.walk(d):
                for fn in files:
                    if Path(fn).suffix.lower() in SUPPORTED_IMG_EXTS:
                        img_files.append(str(Path(root) / fn))
        else:
            fs = filedialog.askopenfilenames(
                title="이미지 파일 선택(다중 가능)",
                filetypes=[("Images", "*.png;*.jpg;*.jpeg;*.webp;*.bmp;*.gif;*.tif;*.tiff")]
            )
            if not fs:
                return
            img_files = list(fs)

        added = 0
        skipped = 0
        for p in img_files:
            try:
                rec = register_image_for_review(
                    out_root=out_dir,
                    image_path=p,
                    origin="manual_add",
                    mode="standalone",
                    parent_doc_id=None,
                    parent_path=None,
                    location={},
                    existing_image_ids=existing_image_ids,
                )
                if rec.get("skipped"):
                    skipped += 1
                else:
                    added += 1
            except Exception:
                continue

        self.log(f"[순수 이미지 추가] added={added}, skipped={skipped}")
        if (added + skipped) > 0:
            self.btn_review.config(state="normal")

    def open_review(self):
        out_dir = self.output_dir.get().strip()
        if not Path(out_dir, "images_queue.jsonl").exists():
            messagebox.showinfo("안내", "images_queue.jsonl이 없습니다.")
            return

        selector = DocumentSelectWindow(self, out_dir)
        self.wait_window(selector)

        if not selector.selected:
           return

        ImageReviewWindow(self, out_dir, selector.selected)



    def export_all_metadata_docs(self):
        out_dir = self.output_dir.get().strip()
        ok, fail, msgs = generate_metadata_docx_for_all(out_dir)
        self.log("\n=== 정제 메타DOCX 생성 결과 ===")
        self.log(f"문서 DOCX: 성공 {ok}, 실패 {fail}")
        for m in msgs[:120]:
            self.log(m)
        if len(msgs) > 120:
            self.log(f"... (총 {len(msgs)}개 중 일부만 표시)")
        messagebox.showinfo("완료", f"정제 메타DOCX 생성 완료\n(output/metadata_docs 폴더 확인)\n- 문서별 DOCX + _IMAGE_DICTIONARY.docx")



    def open_record_builder(self):
        """수기 입력 기반 업무 레코드/보고서 생성 UI"""
        try:
            base_dir = Path.cwd()
            out_dir = Path(self.output_dir.get()).resolve()
            RecordBuilderWindow(self, base_dir=base_dir, out_dir=out_dir)
        except Exception as e:
            messagebox.showerror('오류', f'레코드 창 열기 실패: {e}')

    def open_chunking(self):
        """문서를 선택하여 청킹 처리하는 UI"""
        out_dir = self.output_dir.get().strip()
        if not out_dir or not Path(out_dir).exists():
            messagebox.showerror("오류", "출력 폴더가 설정되지 않았습니다.")
            return
        
        # 문서 선택 창 열기
        selector = DocumentSelectWindow(self, out_dir)
        self.wait_window(selector)
        
        if not selector.selected:
            return
        
        doc_id = selector.selected
        
        # 청킹 설정 창 열기
        settings = ChunkingSettingsWindow(self)
        self.wait_window(settings)
        
        if not settings.confirmed:
            return
        
        # 청킹 처리 실행 (설정값 전달 - IntVar에서 실제 값 가져오기)
        self._run_chunking_for_doc(
            out_dir, 
            doc_id,
            target_chars=settings.target_chars.get(),
            min_chars=settings.min_chars.get(),
            overlap_chars=settings.overlap_chars.get()
        )

    def _run_chunking_for_doc(self, out_root: str, doc_id: str, target_chars: int = 1000, min_chars: int = 250, overlap_chars: int = 120):
        """특정 문서에 대해 청킹 처리 실행"""
        try:
            doc_dir = Path(out_root) / doc_id
            
            # 1) document_facts.json 로드
            facts_path = doc_dir / "document_facts.json"
            if not facts_path.exists():
                messagebox.showerror("오류", f"document_facts.json이 없습니다: {doc_id}")
                return
            
            facts = json.loads(facts_path.read_text(encoding="utf-8"))
            
            # 2) extracted.json에서 텍스트 추출
            extracted_path = doc_dir / "extracted.json"
            if not extracted_path.exists():
                messagebox.showerror("오류", f"extracted.json이 없습니다: {doc_id}")
                return
            
            extracted = json.loads(extracted_path.read_text(encoding="utf-8"))
            text = clean_extracted_text(extracted, max_chars=50000)
            
            if not text.strip():
                messagebox.showwarning("경고", "추출된 텍스트가 없습니다.")
                return
            
            # 3) 청킹 처리 실행
            self.log(f"\n=== 청킹 처리 시작: {doc_id} ===")
            self.log(f"문서 제목: {facts.get('title', '')}")
            self.log(f"텍스트 길이: {len(text)}자")
            self.log(f"청킹 설정: target_chars={target_chars}, min_chars={min_chars}, overlap_chars={overlap_chars}")
            
            try:
                build_chunks_and_reports(
                    out_root=out_root,
                    doc_id=doc_id,
                    facts=facts,
                    text=text,
                    target_chars=target_chars,
                    min_chars=min_chars,
                    overlap_chars=overlap_chars
                )
                self.log(f"✅ 청킹 처리 완료: {doc_id}")
                self.log(f"- output/core/chunks.jsonl에 추가됨")
                self.log(f"- output/{doc_id}/chunks.json 생성됨")
                self.log(f"- output/{doc_id}/analysis_report.json 생성됨")
                messagebox.showinfo("완료", f"청킹 처리 완료\n문서: {facts.get('title', doc_id)}")
            except Exception as e:
                error_msg = f"청킹 처리 실패: {e}"
                self.log(f"❌ {error_msg}\n{traceback.format_exc()}")
                messagebox.showerror("오류", f"청킹 처리 중 오류 발생:\n{e}")
                
        except Exception as e:
            error_msg = f"청킹 처리 준비 실패: {e}"
            self.log(f"❌ {error_msg}\n{traceback.format_exc()}")
            messagebox.showerror("오류", f"청킹 처리 준비 중 오류 발생:\n{e}")

    def open_chunking_viewer(self):
        """청킹 결과를 보기 위한 분석 뷰어 열기"""
        out_dir = self.output_dir.get().strip()
        if not out_dir or not Path(out_dir).exists():
            messagebox.showerror("오류", "출력 폴더가 설정되지 않았습니다.")
            return
        
        # 문서 선택 창 열기
        selector = DocumentSelectWindow(self, out_dir)
        self.wait_window(selector)
        
        if not selector.selected:
            return
        
        doc_id = selector.selected
        
        # 분석 뷰어 열기
        try:
            open_analysis_window(self, out_dir, doc_id)
        except Exception as e:
            self.log(f"⚠️ 분석 뷰어 열기 실패: {e}\n{traceback.format_exc()}")
            messagebox.showerror("오류", f"분석 뷰어를 열 수 없습니다:\n{e}")
    
    def open_tagging_viewer(self):
        """태깅 결과 뷰어 열기"""
        out_dir = self.output_dir.get().strip()
        if not out_dir or not Path(out_dir).exists():
            messagebox.showerror("오류", "출력 폴더가 설정되지 않았습니다.")
            return
        try:
            open_tagging_viewer(self, out_dir)
        except Exception as e:
            self.log(f"⚠️ 태깅 뷰어 열기 실패: {e}\n{traceback.format_exc()}")
            messagebox.showerror("오류", f"태깅 뷰어를 열 수 없습니다:\n{e}")
    
    def open_document_search(self):
        """문서 검색 창 열기"""
        out_dir = self.output_dir.get().strip()
        if not out_dir or not Path(out_dir).exists():
            messagebox.showerror("오류", "출력 폴더가 설정되지 않았습니다.")
            return
        
        try:
            from document_search_ui import DocumentSearchWindow
            DocumentSearchWindow(self, out_dir)
        except Exception as e:
            self.log(f"⚠️ 문서 검색 창 열기 실패: {e}\n{traceback.format_exc()}")
            messagebox.showerror("오류", f"문서 검색 창을 열 수 없습니다:\n{e}")
    
    def open_classification_dashboard(self):
        """문서 분류 대시보드 열기"""
        out_dir = self.output_dir.get().strip()
        if not out_dir or not Path(out_dir).exists():
            messagebox.showerror("오류", "출력 폴더가 설정되지 않았습니다.")
            return
        try:
            from classification_dashboard import open_classification_dashboard
            open_classification_dashboard(self, out_dir)
        except Exception as e:
            self.log(f"⚠️ 문서 분류 대시보드 열기 실패: {e}\n{traceback.format_exc()}")
            messagebox.showerror("오류", f"문서 분류 대시보드를 열 수 없습니다:\n{e}")

if __name__ == "__main__":
    try:
        app = App()
        app.mainloop()
    except Exception as e:
        import traceback
        print("=" * 60)
        print("앱 실행 중 치명적 오류 발생:")
        print("=" * 60)
        traceback.print_exc()
        print("=" * 60)
        input("오류를 확인했습니다. Enter 키를 눌러 종료하세요...")